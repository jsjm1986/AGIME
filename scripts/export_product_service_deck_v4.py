from __future__ import annotations

from pathlib import Path
import math

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

MARGIN_X = 0.45
CARD_GAP = 0.18

BG = RGBColor(0x06, 0x14, 0x2A)
BG_TOP = RGBColor(0x0A, 0x22, 0x43)
CARD_BG = RGBColor(0x0F, 0x28, 0x4A)
CARD_LINE = RGBColor(0x2A, 0x52, 0x7D)
TEXT_MAIN = RGBColor(0xEA, 0xF4, 0xFF)
TEXT_SUB = RGBColor(0x9F, 0xBD, 0xDD)
ACCENT = RGBColor(0x39, 0xD7, 0xFF)
NOTE_BG = RGBColor(0x0E, 0x2E, 0x52)
NOTE_LINE = RGBColor(0x2F, 0xA9, 0xD8)

FONT_FAMILY = "Microsoft YaHei"


def apply_font(run, size: int, bold: bool = False, color: RGBColor = TEXT_MAIN) -> None:
    run.font.name = FONT_FAMILY
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    rpr = run._r.get_or_add_rPr()
    rpr.set(qn("a:ea"), FONT_FAMILY)
    rpr.set(qn("a:cs"), FONT_FAMILY)


def new_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def base_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = BG
    bg_shape.line.fill.background()

    top_bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.23)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = BG_TOP
    top_bar.line.fill.background()
    return slide


def add_header(slide, kicker: str, title: str, subtitle: str, sid: str) -> None:
    kicker_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(MARGIN_X),
        Inches(0.32),
        Inches(3.7),
        Inches(0.33),
    )
    kicker_box.fill.solid()
    kicker_box.fill.fore_color.rgb = BG_TOP
    kicker_box.line.color.rgb = ACCENT
    kicker_box.line.width = Pt(1)
    tf = kicker_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = kicker
    apply_font(run, 10, True, TEXT_MAIN)

    sid_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(11.35),
        Inches(0.32),
        Inches(1.53),
        Inches(0.33),
    )
    sid_box.fill.solid()
    sid_box.fill.fore_color.rgb = BG_TOP
    sid_box.line.color.rgb = CARD_LINE
    sid_box.line.width = Pt(1)
    sid_tf = sid_box.text_frame
    sid_tf.clear()
    sid_p = sid_tf.paragraphs[0]
    sid_p.alignment = PP_ALIGN.CENTER
    sid_run = sid_p.add_run()
    sid_run.text = sid
    apply_font(sid_run, 10, True, TEXT_SUB)

    title_box = slide.shapes.add_textbox(Inches(MARGIN_X), Inches(0.72), Inches(12.3), Inches(0.88))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    apply_font(run, 29, True, TEXT_MAIN)

    subtitle_box = slide.shapes.add_textbox(Inches(MARGIN_X), Inches(1.52), Inches(12.3), Inches(0.38))
    sub_tf = subtitle_box.text_frame
    sub_tf.clear()
    sub_p = sub_tf.paragraphs[0]
    sub_p.alignment = PP_ALIGN.LEFT
    sub_run = sub_p.add_run()
    sub_run.text = subtitle
    apply_font(sub_run, 12, False, TEXT_SUB)


def add_metric_row(slide, top_y: float, metrics: list[tuple[str, str]]) -> None:
    total_w = 13.333 - (MARGIN_X * 2)
    cols = max(1, len(metrics))
    card_w = (total_w - CARD_GAP * (cols - 1)) / cols
    card_h = 1.03

    for i, (num, text) in enumerate(metrics):
        x = MARGIN_X + i * (card_w + CARD_GAP)
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(top_y),
            Inches(card_w),
            Inches(card_h),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.color.rgb = CARD_LINE
        shape.line.width = Pt(1)

        tf = shape.text_frame
        tf.clear()
        tf.margin_left = Inches(0.14)
        tf.margin_right = Inches(0.14)
        tf.margin_top = Inches(0.06)
        tf.margin_bottom = Inches(0.06)

        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.LEFT
        r1 = p1.add_run()
        r1.text = num
        apply_font(r1, 23, True, ACCENT)
        p1.space_after = Pt(0)

        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = text
        apply_font(r2, 10, False, TEXT_SUB)


def add_card(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    lines: list[str],
    body_size: int = 11,
) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = CARD_LINE
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.16)
    tf.margin_right = Inches(0.16)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.08)

    p_title = tf.paragraphs[0]
    p_title.alignment = PP_ALIGN.LEFT
    r_title = p_title.add_run()
    r_title.text = title
    apply_font(r_title, 15, True, TEXT_MAIN)
    p_title.space_after = Pt(2)

    for line in lines:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = f"• {line}"
        apply_font(r, body_size, False, TEXT_SUB)
        p.space_after = Pt(1)


def add_cards_block(
    slide,
    top_y: float,
    block_h: float,
    cards: list[dict],
    cols: int,
    body_size: int = 11,
) -> None:
    total_w = 13.333 - (MARGIN_X * 2)
    card_w = (total_w - CARD_GAP * (cols - 1)) / cols
    rows = math.ceil(len(cards) / cols)
    row_h = (block_h - CARD_GAP * (rows - 1)) / rows

    for idx, card in enumerate(cards):
        row = idx // cols
        col = idx % cols
        x = MARGIN_X + col * (card_w + CARD_GAP)
        y = top_y + row * (row_h + CARD_GAP)
        add_card(slide, x, y, card_w, row_h, card["title"], card["lines"], body_size=body_size)


def add_note(slide, y: float, text: str, h: float = 0.62) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(MARGIN_X),
        Inches(y),
        Inches(12.43),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = NOTE_BG
    shape.line.color.rgb = NOTE_LINE
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.14)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    apply_font(run, 11, False, TEXT_MAIN)


def add_footer(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(MARGIN_X), Inches(7.15), Inches(12.43), Inches(0.24))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    apply_font(run, 9, False, TEXT_SUB)


def build_product_service_slides(prs: Presentation) -> None:
    # 01
    s = base_slide(prs)
    add_header(
        s,
        "PRODUCT DECK · 2026-02-24",
        "AGIME 产品与服务说明书：把 AI 工具变成团队执行系统",
        "目标：让第一次接触 AGIME 的客户 10 分钟看懂“产品是什么、能做什么、我们卖什么”。",
        "SLIDE 01 / 16",
    )
    add_metric_row(
        s,
        1.86,
        [
            ("开源免费", "软件本体永久免费"),
            ("本地优先", "个人数据默认可控"),
            ("Team Server", "团队共享与治理执行"),
            ("4 档套餐", "199/4999/8999/29999"),
        ],
    )
    add_cards_block(
        s,
        3.05,
        2.75,
        [
            {
                "title": "这份 PPT 回答 3 个问题",
                "lines": [
                    "AGIME 到底是什么产品？",
                    "它能给个人和团队带来什么结果？",
                    "当前我们具体在卖哪些套餐与服务？",
                ],
            },
            {
                "title": "一句话定位",
                "lines": [
                    "不只是一个 AI 工具，而是一套协作新范式",
                    "把个人能力、团队能力和外部协作放进同一执行系统",
                    "让“个人会做”升级为“团队都会做”",
                ],
            },
        ],
        cols=2,
        body_size=11,
    )
    add_note(s, 5.92, "叙事原则：先讲产品价值，再讲套餐服务；先讲能做什么，再讲怎么买。")
    add_footer(s, "下一页：AGIME 产品名与核心理念。")

    # 02
    s = base_slide(prs)
    add_header(
        s,
        "PRODUCT DEFINITION",
        "AGIME 是什么：Human × AI × Team × Ecosystem",
        "官网核心定义：Human 负责决策，AI 负责执行，Team 负责复用治理，Ecosystem 负责协同交付。",
        "SLIDE 02 / 16",
    )
    add_cards_block(
        s,
        1.9,
        4.2,
        [
            {
                "title": "产品名说明",
                "lines": [
                    "AGIME 可以理解为 AI + Me 的组织化延展",
                    "不是替代人，而是放大人的执行效率",
                    "重点不是聊天，而是可执行与可复用",
                ],
            },
            {
                "title": "产品边界",
                "lines": [
                    "不是只做对话助手",
                    "不是只做个人插件合集",
                    "而是个人端 + 团队端 + 对外协作端的一体化系统",
                ],
            },
            {
                "title": "协作链路（原文）",
                "lines": [
                    "个人沉淀能力 -> 团队共享复用",
                    "Agent 云端执行 -> 对外协作交付",
                    "最终形成组织级执行系统",
                ],
            },
            {
                "title": "你可以怎么理解它",
                "lines": [
                    "个人版：你的本地 AI 执行助手",
                    "Team Server：团队版 AI 工作中台",
                    "Portal：对外统一发布与协作窗口",
                ],
            },
        ],
        cols=2,
        body_size=10,
    )
    add_footer(s, "下一页：它解决的真实问题。")

    # 03
    s = base_slide(prs)
    add_header(
        s,
        "PAIN POINT",
        "为什么需要 AGIME：旧工作方式存在结构性瓶颈",
        "很多团队不是没有 AI，而是 AI 无法稳定接入日常执行流程。",
        "SLIDE 03 / 16",
    )
    add_cards_block(
        s,
        1.9,
        4.9,
        [
            {"title": "隐私与效率难兼得", "lines": ["想提效但担心核心数据外流", "流程结果经常回到人工收尾"]},
            {"title": "工具分散成本堆叠", "lines": ["多人多工具、多配置、多订阅", "投入增加但协同效率不升"]},
            {"title": "会建议，不会执行", "lines": ["AI 会给步骤，但执行仍靠手工", "重复劳动长期存在"]},
            {"title": "环境受限能力断层", "lines": ["内网、跨系统、跨地域场景复杂", "工具无法稳定持续执行"]},
            {"title": "经验难沉淀", "lines": ["高手方法停留在个人对话里", "新人只能重复试错"]},
            {"title": "团队复用困难", "lines": ["共享后缺权限与治理边界", "流程越做越乱"]},
        ],
        cols=3,
        body_size=9,
    )
    add_footer(s, "下一页：AGIME 的产品结构。")

    # 04
    s = base_slide(prs)
    add_header(
        s,
        "PRODUCT ARCH",
        "产品结构：个人层 + 团队层 + 生态层",
        "核心不是多功能叠加，而是把“执行、复用、治理”放进一条链路。",
        "SLIDE 04 / 16",
    )
    add_cards_block(
        s,
        1.9,
        3.7,
        [
            {
                "title": "个人协作层（本地）",
                "lines": [
                    "文件处理、自动化执行、任务沉淀",
                    "本地优先，数据可控",
                    "适合个人和小团队快速起步",
                ],
            },
            {
                "title": "团队协作层（Team Server）",
                "lines": [
                    "共享 Skills/MCP/文档知识",
                    "角色权限、审核、审计治理",
                    "云端 Team Agent 持续执行任务",
                ],
            },
            {
                "title": "生态协同层（Portal）",
                "lines": [
                    "对外发布统一资料与答疑",
                    "连接客户、供应商、合作方",
                    "把内部能力标准化输出",
                ],
            },
        ],
        cols=3,
        body_size=10,
    )
    add_metric_row(
        s,
        5.83,
        [
            ("执行", "不是只回答，而是可落地做事"),
            ("复用", "一人沉淀，全员可用"),
            ("治理", "权限、审计、边界清晰"),
        ],
    )
    add_note(s, 6.95, "关系说明：Team 不是替代个人版，而是在个人能力之上增加共享、治理和云端执行。", h=0.4)

    # 05
    s = base_slide(prs)
    add_header(
        s,
        "PERSONAL CAPABILITY",
        "个人层能做什么：从文档到自动化的日常执行",
        "重点是“可执行”和“可复用”，而不是一次性对话。",
        "SLIDE 05 / 16",
    )
    add_cards_block(
        s,
        1.9,
        4.9,
        [
            {"title": "文档批处理", "lines": ["读取 PDF/Word/Excel", "提取结构化信息", "自动汇总成报告"]},
            {"title": "桌面自动化", "lines": ["模拟点击和键盘输入", "跨系统重复操作自动化", "减少手工复制粘贴"]},
            {"title": "规范沉淀", "lines": ["一次对话后提取可复用流程", "下次同类任务直接复用", "降低重复沟通"]},
            {"title": "定时执行", "lines": ["支持定时定量任务", "无需人工持续盯盘", "降低遗漏风险"]},
            {"title": "网页调研", "lines": ["自动浏览并采集信息", "按模板整理结果", "用于竞品与行业研究"]},
            {"title": "数据分析", "lines": ["辅助分析数据与异常", "生成图表与摘要", "支持决策汇报材料"]},
        ],
        cols=3,
        body_size=9,
    )
    add_footer(s, "下一页：Team Server 的团队价值。")

    # 06
    s = base_slide(prs)
    add_header(
        s,
        "TEAM SERVER",
        "团队层到底有什么用：把个人效率升级为组织效率",
        "任务交给云端 Team Agent 持续执行，团队共享技能、工具、知识。",
        "SLIDE 06 / 16",
    )
    add_cards_block(
        s,
        1.9,
        4.05,
        [
            {
                "title": "共享执行中枢",
                "lines": ["云端持续执行，不占个人电脑", "自动调用共享 Skills/MCP/知识库", "结果沉淀为团队资产"],
            },
            {
                "title": "治理能力",
                "lines": ["角色分工、权限管理、审核流程", "关键动作可追踪可审计", "避免“共享即失控”"],
            },
            {
                "title": "协作能力",
                "lines": ["新人快速上手已有流程", "跨地域团队统一版本", "Portal 对外统一窗口"],
            },
            {
                "title": "部署方式",
                "lines": ["同办公室：LAN 协作", "跨地域：云端 Team Server", "团队正式协作建议 MongoDB 模式"],
            },
        ],
        cols=2,
        body_size=10,
    )
    add_metric_row(
        s,
        6.16,
        [
            ("一人沉淀", "方法不再锁在个人会话"),
            ("全员复用", "新人加入即可调用标准流程"),
            ("持续执行", "重复工作交给 Team Agent"),
        ],
    )
    add_footer(s, "下一页：典型客户怎么用。")

    # 07
    s = base_slide(prs)
    add_header(
        s,
        "WHO USES AGIME",
        "谁最适合用 AGIME：从个人到企业的四类用户",
        "不是限定单一角色，而是覆盖执行者、管理者和交付团队。",
        "SLIDE 07 / 16",
    )
    add_cards_block(
        s,
        1.9,
        3.9,
        [
            {"title": "个人执行者", "lines": ["文档处理与自动化", "沉淀个人流程模板", "减少重复事务"]},
            {"title": "团队负责人", "lines": ["统一流程与权限", "推动跨岗位复用", "降低协作扯皮成本"]},
            {"title": "企业管理层", "lines": ["关注可控与可审计", "从试点走向规模落地", "沉淀组织能力资产"]},
            {"title": "交付与技术伙伴", "lines": ["基于 MCP/二开做行业方案", "复用标准交付模板", "缩短项目周期"]},
        ],
        cols=2,
        body_size=10,
    )
    add_metric_row(
        s,
        6.0,
        [
            ("起步快", "先用个人版验证高频任务"),
            ("升级顺", "再接 Team Server 做团队复用"),
            ("扩展稳", "最后进入私有化与定制交付"),
        ],
    )
    add_footer(s, "下一页：从体验到上线的落地路径。")

    # 08
    s = base_slide(prs)
    add_header(
        s,
        "ONBOARDING PATH",
        "落地路径：7 天看效果，30 天出样板，90 天可复制",
        "先解决实际问题，再决定规模化投入。",
        "SLIDE 08 / 16",
    )
    add_cards_block(
        s,
        1.9,
        4.8,
        [
            {
                "title": "第 1 阶段（Day 1-7）",
                "lines": ["安装体验 + 选 1-2 个高频任务", "跑通文档/自动化小闭环", "确认基础可用性"],
            },
            {
                "title": "第 2 阶段（Day 8-30）",
                "lines": ["沉淀标准模板和 SOP", "让 2-3 个成员复用同流程", "产出首个可汇报样板"],
            },
            {
                "title": "第 3 阶段（Day 31-90）",
                "lines": ["接入 Team Server 做持续执行", "建立权限与审计边界", "形成跨部门可复制路径"],
            },
            {
                "title": "决策闸门",
                "lines": ["如果结果稳定：进入立项与定制", "如果场景不匹配：保持轻量使用", "避免一次性重投入"],
            },
            {
                "title": "成功标准（非财务）",
                "lines": ["重复任务自动化率提升", "团队复用率提升", "流程执行稳定性提升"],
            },
            {
                "title": "常见失败原因",
                "lines": ["范围过大，没从高频任务开始", "缺统一验收口径", "没有把流程沉淀成模板"],
            },
        ],
        cols=3,
        body_size=9,
    )
    add_footer(s, "下一页：对比实验室怎么做。")

    # 09
    s = base_slide(prs)
    add_header(
        s,
        "COMPARISON LAB",
        "对比实验室：我们怎么做横向与纵向对比",
        "仅针对“用户向桌面 Agent”，避免和编程 Agent 混口径。",
        "SLIDE 09 / 16",
    )
    add_metric_row(
        s,
        1.86,
        [
            ("11", "用户向桌面产品样本"),
            ("4", "关键评估轴线"),
            ("25+", "官方资料交叉验证"),
            ("双编号", "S 结论 / D 文档可审计"),
        ],
    )
    add_cards_block(
        s,
        3.05,
        3.32,
        [
            {
                "title": "横向矩阵关注点",
                "lines": ["本地离线能力", "文件/系统执行能力", "自动化与工作流", "团队治理与私有部署"],
            },
            {
                "title": "纵向评分三主轴",
                "lines": ["执行深度", "隐私与部署控制", "团队治理与可扩展"],
            },
            {
                "title": "评分口径",
                "lines": ["基于官方公开能力映射", "用于快速筛选，不替代 PoC", "支持证据追溯复核"],
            },
            {
                "title": "典型样本（精选）",
                "lines": ["AGIME / LobsterAI / QoderWork / OpenClaw", "ChatGPT Desktop / Claude Desktop", "Perplexity / Jan / Copilot 等"],
            },
        ],
        cols=2,
        body_size=10,
    )
    add_footer(s, "下一页：对比结论精选。")

    # 10
    s = base_slide(prs)
    add_header(
        s,
        "COMPARISON RESULT",
        "对比结论精选：AGIME 强项在“执行 + 可控 + 可治理”",
        "这里不追求全表堆砌，只放最影响决策的结论。",
        "SLIDE 10 / 16",
    )
    add_cards_block(
        s,
        1.9,
        4.24,
        [
            {
                "title": "执行深度维度",
                "lines": ["AGIME 在文件/系统执行与流程复用上更完整", "适合从“会回答”升级到“会做事”场景", "对比中同类产品常在执行闭环上有短板"],
            },
            {
                "title": "隐私与部署维度",
                "lines": ["本地优先 + 私有化路径，适配数据边界要求", "比纯云助手更适合敏感业务场景", "个人与企业可用同一能力体系演进"],
            },
            {
                "title": "团队治理维度",
                "lines": ["Team Server 支持共享、权限、审计", "更容易从个人提效升级为组织提效", "降低跨部门协作失控风险"],
            },
            {
                "title": "决策建议",
                "lines": ["先矩阵筛选 2-3 个候选", "再做 2 周 PoC 看真实执行效果", "优先选择能沉淀组织能力的方案"],
            },
        ],
        cols=2,
        body_size=10,
    )
    add_note(s, 6.28, "判断标准：不是谁“更会说”，而是谁“更能稳定交付结果”。")
    add_footer(s, "下一页：为什么客户最终会选 AGIME。")

    # 11
    s = base_slide(prs)
    add_header(
        s,
        "WHY AGIME",
        "为什么选 AGIME：四个可感知差异",
        "从产品体验到组织落地，强调可执行、可复用、可治理。",
        "SLIDE 11 / 16",
    )
    add_cards_block(
        s,
        1.9,
        4.8,
        [
            {"title": "差异 1：执行不是演示", "lines": ["文件、系统、自动化可真实落地", "不止给建议，能替你执行"]},
            {"title": "差异 2：本地优先可控", "lines": ["个人可控，企业可管", "隐私与效率不再二选一"]},
            {"title": "差异 3：团队能力可沉淀", "lines": ["一人沉淀，全员复用", "经验变资产，不靠单点高手"]},
            {"title": "差异 4：服务边界清晰", "lines": ["标准包先验证，再进入定制", "套餐和交付边界清楚，决策更省心"]},
            {"title": "结果导向", "lines": ["先证明能做，再谈规模化", "减少试错成本与决策风险"]},
            {"title": "可持续性", "lines": ["开源入口降低门槛", "服务产品化保障长期交付质量"]},
        ],
        cols=3,
        body_size=9,
    )
    add_footer(s, "下一页：我们到底在卖什么。")

    # 12
    s = base_slide(prs)
    add_header(
        s,
        "OFFER CATALOG",
        "我们卖什么：产品免费 + 服务收费",
        "核心原则：开源代码免费，收费项对应服务与交付，不卖封闭许可证。",
        "SLIDE 12 / 16",
    )
    add_cards_block(
        s,
        1.9,
        4.1,
        [
            {
                "title": "免费部分（产品）",
                "lines": ["AGIME 开源软件本体（Apache-2.0）", "个人本地执行能力", "基础使用文档与社区支持"],
            },
            {
                "title": "收费部分（标准服务包）",
                "lines": ["199 诊断启动包", "4999 部署验证包", "8999 业务试点包", "29,999 定制立项包"],
            },
            {
                "title": "收费部分（企业服务）",
                "lines": ["私有化部署与去 Logo 授权", "MCP Skills 定向开发", "AGIME 二开与系统集成"],
            },
            {
                "title": "我们不卖什么",
                "lines": ["不卖封闭源码许可", "不把 4999/8999 包装成定制开发", "不做无边界“无限人力陪跑”"],
            },
        ],
        cols=2,
        body_size=10,
    )
    add_metric_row(
        s,
        6.18,
        [
            ("先标准包", "先验证可行性与匹配度"),
            ("再立项包", "明确范围、边界、里程碑"),
            ("后定制开发", "按里程碑交付与验收"),
        ],
    )
    add_footer(s, "下一页：199 与 4999 套餐明细。")

    # 13
    s = base_slide(prs)
    add_header(
        s,
        "PACKAGE A/B",
        "套餐明细（上）：199 诊断包 & 4999 验证包",
        "这两档目标是“低成本验证与筛选”，不是定制开发。",
        "SLIDE 13 / 16",
    )
    add_cards_block(
        s,
        1.9,
        4.86,
        [
            {
                "title": "199 诊断启动包｜¥199",
                "lines": [
                    "标准化入门教程库（参考价值 499）",
                    "$99 Token 启动包（按购买日汇率折算）",
                    "自动化诊断问卷与评分 + 标准任务模板",
                    "输出：系统生成诊断报告，适合快速判断是否匹配",
                ],
            },
            {
                "title": "4999 部署验证包｜¥4,999",
                "lines": [
                    "私有化部署脚本与指南（参考价值 3,999）",
                    "去 Logo 商用授权（参考价值 6,999）",
                    "$299 Token 验证包 + 标准验收清单文档",
                    "输出：完成标准部署验证，不含定制开发",
                ],
            },
            {
                "title": "这两档适合谁",
                "lines": [
                    "想低风险试错、先看结果的团队",
                    "需要快速判断“能不能落地”的负责人",
                    "不希望一开始就进入重定制项目",
                ],
            },
            {
                "title": "边界提醒",
                "lines": [
                    "199/4999 为一次性标准包",
                    "不默认包含深度定制开发",
                    "若需定制，建议走 29,999 立项包",
                ],
            },
        ],
        cols=2,
        body_size=9,
    )
    add_footer(s, "下一页：8999 与 29,999 套餐明细。")

    # 14
    s = base_slide(prs)
    add_header(
        s,
        "PACKAGE C/D",
        "套餐明细（下）：8999 试点包 & 29,999 立项包",
        "目标是“从可用走向可落地”，并为定制交付建立清晰边界。",
        "SLIDE 14 / 16",
    )
    add_cards_block(
        s,
        1.9,
        4.86,
        [
            {
                "title": "8999 业务试点包｜¥8,999",
                "lines": [
                    "标准化试点模板库（参考价值 5,999）",
                    "团队自助教程库 + $499 Token 试跑包",
                    "自动化复盘报表与需求模板",
                    "输出：试点可行性结论，不含定制开发",
                ],
            },
            {
                "title": "29,999 定制立项包｜¥29,999",
                "lines": [
                    "需求采集模板 + 技术架构蓝图",
                    "$999 Token 研发验证包 + 里程碑报价模板",
                    "签约后可 100% 抵扣定制合同首款",
                    "输出：项目范围/边界/排期/报价的立项依据",
                ],
            },
            {
                "title": "为什么先立项再开发",
                "lines": [
                    "先把范围讲清楚，减少返工",
                    "先把验收讲清楚，避免扯皮",
                    "先把里程碑讲清楚，便于管理进度",
                ],
            },
            {
                "title": "交付方式",
                "lines": [
                    "立项后按里程碑开发交付",
                    "按阶段验收和回款",
                    "从“模糊需求”变“可执行合同”",
                ],
            },
        ],
        cols=2,
        body_size=9,
    )
    add_footer(s, "下一页：企业级扩展服务。")

    # 15
    s = base_slide(prs)
    add_header(
        s,
        "ENTERPRISE SERVICE",
        "企业级扩展服务：我们还能提供什么",
        "当标准包验证通过后，可进入专项服务与定制交付。",
        "SLIDE 15 / 16",
    )
    add_cards_block(
        s,
        1.9,
        4.86,
        [
            {
                "title": "私有化部署与授权",
                "lines": [
                    "私有化环境部署与配置",
                    "去 Logo 商用授权",
                    "适配企业数据边界与安全要求",
                ],
            },
            {
                "title": "MCP Skills 定向开发",
                "lines": [
                    "按业务流程定制技能与工具链",
                    "沉淀行业模板，提升复用效率",
                    "适合中大型流程标准化改造",
                ],
            },
            {
                "title": "AGIME 二次开发",
                "lines": [
                    "面向企业场景做功能扩展",
                    "对接内部系统与业务流程",
                    "以里程碑方式推进交付",
                ],
            },
            {
                "title": "培训与上线支持",
                "lines": [
                    "角色化培训与标准教材",
                    "上线初期流程校准与复盘",
                    "帮助团队快速形成可复制实践",
                ],
            },
            {
                "title": "系统集成与治理",
                "lines": [
                    "权限、审计、流程边界落地",
                    "统一工单与交付文档体系",
                    "支撑跨部门持续协作",
                ],
            },
            {
                "title": "你拿到的结果",
                "lines": [
                    "不是一次性项目，而是持续可复用能力",
                    "组织可治理、可审计、可复制",
                    "后续扩展成本更可控",
                ],
            },
        ],
        cols=3,
        body_size=9,
    )
    add_footer(s, "下一页：常见问题与下一步行动。")

    # 16
    s = base_slide(prs)
    add_header(
        s,
        "FAQ · NEXT STEP",
        "常见问题 + 下一步：先体验，再选包，再扩展",
        "这页给到决策者一个简单、可执行的启动路径。",
        "SLIDE 16 / 16",
    )
    add_cards_block(
        s,
        1.9,
        3.9,
        [
            {
                "title": "FAQ 1：AGIME 真免费吗？",
                "lines": ["软件本体永久免费开源", "收费项是服务、授权、培训与定制交付", "可先免费体验再决定购买服务"],
            },
            {
                "title": "FAQ 2：数据安全吗？",
                "lines": ["本地优先，文件可不出本机", "可选私有化部署路径", "团队版支持权限和审计治理"],
            },
            {
                "title": "FAQ 3：Team 会替代个人版吗？",
                "lines": ["不会", "Team 是在个人能力上增加共享/治理/云端执行", "适合多人协作与跨部门复用"],
            },
            {
                "title": "建议启动路径",
                "lines": ["第 1 步：下载个人版跑一个真实任务", "第 2 步：选择 199/4999/8999 做验证", "第 3 步：匹配后进入 29,999 立项与企业服务"],
            },
        ],
        cols=2,
        body_size=10,
    )
    add_note(s, 5.96, "资料来源：官网 index.html（产品与 Team 叙事）｜product-comparison.html（对比实验）｜business-channel.html（套餐与服务）。")
    add_metric_row(
        s,
        6.55,
        [
            ("产品优先", "先让用户看懂产品再谈扩展"),
            ("边界清晰", "每档套餐输出与不包含项明确"),
            ("可执行", "可以直接用于客户介绍与销售沟通"),
        ],
    )


def main() -> None:
    prs = new_presentation()
    build_product_service_slides(prs)

    out_dir = Path("wwwweb")
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / "AGIME_product_service_deck_v4_16slides_2026-02-24.pptx"
    prs.save(str(out_path))
    print(f"Generated: {out_path}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
