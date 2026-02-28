import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
from pptx import Presentation

p = Presentation('data/workspaces/698616a1980c003c66f6421e/missions/4f46f0b7-52aa-4af6-9b40-c64646c75ab1/output/信任的杠杆_中医AI平台_增强版.pptx')
for i, s in enumerate(p.slides, 1):
    texts = [sh.text_frame.text.strip() for sh in s.shapes if sh.has_text_frame and sh.text_frame.text.strip()]
    print(f'\n=== 第{i}页 ===')
    for t in texts[:6]:
        print(t[:80])
