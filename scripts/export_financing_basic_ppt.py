from __future__ import annotations

from pathlib import Path
import math

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

MARGIN_X = 0.45
CARD_GAP = 0.18

BG = RGBColor(0x06, 0x14, 0x2A)
BG_TOP = RGBColor(0x0A, 0x22, 0x43)
CARD_BG = RGBColor(0x0F, 0x28, 0x4A)
CARD_LINE = RGBColor(0x2A, 0x52, 0x7D)
TEXT_MAIN = RGBColor(0xEA, 0xF4, 0xFF)
TEXT_SUB = RGBColor(0x9F, 0xBD, 0xDD)
ACCENT = RGBColor(0x39, 0xD7, 0xFF)
NOTE_BG = RGBColor(0x0E, 0x2E, 0x52)
NOTE_LINE = RGBColor(0x2F, 0xA9, 0xD8)

FONT_FAMILY = "Microsoft YaHei"


def apply_font(run, size: int, bold: bool = False, color: RGBColor = TEXT_MAIN) -> None:
    run.font.name = FONT_FAMILY
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    # Ensure East Asian font displays correctly.
    rpr = run._r.get_or_add_rPr()
    rpr.set(qn("a:ea"), FONT_FAMILY)
    rpr.set(qn("a:cs"), FONT_FAMILY)


def new_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def base_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = BG
    bg_shape.line.fill.background()

    top_bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.23)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = BG_TOP
    top_bar.line.fill.background()

    return slide


def add_header(slide, kicker: str, title: str, subtitle: str, sid: str) -> None:
    kicker_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(MARGIN_X),
        Inches(0.32),
        Inches(3.5),
        Inches(0.33),
    )
    kicker_box.fill.solid()
    kicker_box.fill.fore_color.rgb = BG_TOP
    kicker_box.line.color.rgb = ACCENT
    kicker_box.line.width = Pt(1)
    tf = kicker_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = kicker
    apply_font(run, 10, True, TEXT_MAIN)

    sid_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(11.35),
        Inches(0.32),
        Inches(1.53),
        Inches(0.33),
    )
    sid_box.fill.solid()
    sid_box.fill.fore_color.rgb = BG_TOP
    sid_box.line.color.rgb = CARD_LINE
    sid_box.line.width = Pt(1)
    sid_tf = sid_box.text_frame
    sid_tf.clear()
    sid_p = sid_tf.paragraphs[0]
    sid_p.alignment = PP_ALIGN.CENTER
    sid_run = sid_p.add_run()
    sid_run.text = sid
    apply_font(sid_run, 10, True, TEXT_SUB)

    title_box = slide.shapes.add_textbox(Inches(MARGIN_X), Inches(0.72), Inches(12.3), Inches(0.88))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    apply_font(run, 30, True, TEXT_MAIN)

    subtitle_box = slide.shapes.add_textbox(Inches(MARGIN_X), Inches(1.52), Inches(12.3), Inches(0.38))
    sub_tf = subtitle_box.text_frame
    sub_tf.clear()
    sub_p = sub_tf.paragraphs[0]
    sub_p.alignment = PP_ALIGN.LEFT
    sub_run = sub_p.add_run()
    sub_run.text = subtitle
    apply_font(sub_run, 12, False, TEXT_SUB)


def add_metric_row(slide, top_y: float, metrics: list[tuple[str, str]]) -> None:
    total_w = 13.333 - (MARGIN_X * 2)
    cols = max(1, len(metrics))
    card_w = (total_w - CARD_GAP * (cols - 1)) / cols
    card_h = 1.03

    for i, (num, text) in enumerate(metrics):
        x = MARGIN_X + i * (card_w + CARD_GAP)
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(top_y),
            Inches(card_w),
            Inches(card_h),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.color.rgb = CARD_LINE
        shape.line.width = Pt(1)

        tf = shape.text_frame
        tf.clear()
        tf.margin_left = Inches(0.14)
        tf.margin_right = Inches(0.14)
        tf.margin_top = Inches(0.06)
        tf.margin_bottom = Inches(0.06)

        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.LEFT
        r1 = p1.add_run()
        r1.text = num
        apply_font(r1, 23, True, ACCENT)
        p1.space_after = Pt(0)

        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = text
        apply_font(r2, 10, False, TEXT_SUB)


def add_card(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    lines: list[str],
    body_size: int = 11,
) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.color.rgb = CARD_LINE
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.16)
    tf.margin_right = Inches(0.16)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.08)

    p_title = tf.paragraphs[0]
    p_title.alignment = PP_ALIGN.LEFT
    r_title = p_title.add_run()
    r_title.text = title
    apply_font(r_title, 15, True, TEXT_MAIN)
    p_title.space_after = Pt(2)

    for line in lines:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = f"• {line}"
        apply_font(r, body_size, False, TEXT_SUB)
        p.space_after = Pt(1)


def add_cards_block(
    slide,
    top_y: float,
    block_h: float,
    cards: list[dict],
    cols: int,
    body_size: int = 11,
) -> None:
    total_w = 13.333 - (MARGIN_X * 2)
    card_w = (total_w - CARD_GAP * (cols - 1)) / cols
    rows = math.ceil(len(cards) / cols)
    row_h = (block_h - CARD_GAP * (rows - 1)) / rows

    for idx, card in enumerate(cards):
        row = idx // cols
        col = idx % cols
        x = MARGIN_X + col * (card_w + CARD_GAP)
        y = top_y + row * (row_h + CARD_GAP)
        add_card(slide, x, y, card_w, row_h, card["title"], card["lines"], body_size=body_size)


def add_note(slide, y: float, text: str, h: float = 0.62) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(MARGIN_X),
        Inches(y),
        Inches(12.43),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = NOTE_BG
    shape.line.color.rgb = NOTE_LINE
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.14)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    apply_font(run, 11, False, TEXT_MAIN)


def add_footer(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(MARGIN_X), Inches(7.15), Inches(12.43), Inches(0.24))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    apply_font(run, 9, False, TEXT_SUB)


def build_basic_slides(prs: Presentation) -> None:
    # B1
    s = base_slide(prs)
    add_header(
        s,
        "NON-TECH DECISION · 2026-02-24",
        "AGIME：把 AI 焦虑变成可复制收益",
        "面向非技术背景的投资决策场景：市场在焦虑，客户在试错，我们提供可交付、可复制、可放大的落地方案。",
        "SLIDE 01 / 10",
    )
    add_metric_row(
        s,
        1.86,
        [
            ("$2.52T", "2026 全球 AI 支出仍高速增长"),
            ("88%", "组织已开始使用 AI，但效果不稳定"),
            ("3 条", "收入主引擎：服务包、渠道收入、定制交付"),
            ("6 步", "闭环：诊断->模板->执行->复盘->SOP->扩散"),
        ],
    )
    add_cards_block(
        s,
        3.03,
        2.65,
        [
            {
                "title": "核心决策问题（3条）",
                "lines": [
                    "这是不是刚需，还是短期噱头？",
                    "收入是否可复制，不依赖单一大客户？",
                    "团队优势是否会被快速替代？",
                ],
            },
            {
                "title": "我们的直接回答",
                "lines": [
                    "需求真实且持续：客户焦虑的是“不会落地”",
                    "收入结构是组合型：服务包 + 渠道 + 定制",
                    "壁垒来自交付系统化，而不是单点功能",
                ],
            },
        ],
        cols=2,
        body_size=11,
    )
    add_note(s, 5.78, "一句话：我们不靠“讲模型故事”赚钱，而是靠“把客户效率稳定做出来”赚钱。")
    add_footer(s, "下一页：先讲清首页口号与产品名到底意味着什么。")

    # New: Homepage slogan and product naming narrative
    s = base_slide(prs)
    add_header(
        s,
        "HOMEPAGE NARRATIVE",
        "首页口号与产品名说明：AGIME 是协作系统，不是聊天工具",
        "来自官网首屏的核心定义：从 Human × AI 到 Team × Ecosystem 的可执行协作链路。",
        "SLIDE 02 / 10",
    )
    add_cards_block(
        s,
        1.9,
        3.55,
        [
            {
                "title": "口号（官网原文）",
                "lines": [
                    "不只是一个 AI 工具，而是一套协作新范式",
                    "Human × AI × Team × Ecosystem",
                    "把个人能力、团队能力和外部协作放进同一套可执行系统",
                ],
            },
            {
                "title": "产品名说明（AGIME）",
                "lines": [
                    "可理解为 AI + Me 的组织化延展",
                    "Human 负责决策，AI 负责执行",
                    "Team 负责复用与治理，Ecosystem 负责跨组织协同与交付",
                ],
            },
            {
                "title": "协作链路（官网原文）",
                "lines": [
                    "个人沉淀能力 -> 团队共享复用",
                    "Team Agent 云端执行 -> 对外协作交付",
                    "把“个人会做”变成“团队都会做”",
                ],
            },
            {
                "title": "双层能力结构",
                "lines": [
                    "个人协作层：文件处理、自动化执行、本地优先",
                    "团队协作层（Team Server）：共享技能、权限治理、审计追踪",
                    "能力沉淀后可跨部门、跨组织复制",
                ],
            },
        ],
        cols=2,
        body_size=10,
    )
    add_metric_row(
        s,
        5.62,
        [
            ("永久免费开源", "降低试错成本，先验证再规模化"),
            ("本地优先可控", "个人与企业均可守住数据边界"),
            ("支持团队私有化", "从个人效率升级为组织协同效率"),
        ],
    )
    add_note(s, 6.88, "这页解决投资人第一问：AGIME 不是“又一个模型入口”，而是“把执行与治理合并”的协作基础设施。", h=0.4)
    add_footer(s, "下一页：市场机会与需求缺口。")

    # B2
    s = base_slide(prs)
    add_header(
        s,
        "MARKET GAP",
        "市场机会：大家都在买 AI，但多数企业不会系统提效",
        "买方最痛的不是“没有工具”，而是“花了钱却没有稳定产出”。",
        "SLIDE 03 / 10",
    )
    add_cards_block(
        s,
        1.9,
        2.46,
        [
            {
                "title": "现实痛点",
                "lines": ["团队会用 AI，但结果波动大", "项目试点多，规模化少", "成本投入后很难复盘 ROI"],
            },
            {
                "title": "市场错位",
                "lines": [
                    "市场供给偏“工具”，需求偏“结果”",
                    "客户缺的是落地方法，不是更多功能",
                    "企业担心安全与治理风险",
                ],
            },
            {
                "title": "直接机会",
                "lines": [
                    "谁能稳定交付结果，谁就能拿预算",
                    "谁能复制交付，谁就能做规模",
                    "谁能沉淀模板，谁就有长期壁垒",
                ],
            },
        ],
        cols=3,
        body_size=11,
    )
    add_cards_block(
        s,
        4.52,
        1.64,
        [
            {
                "title": "我们服务的客户",
                "lines": ["探索 AI 转型但“方向不清、团队不会、成本不可控”的企业与团队"],
            },
            {
                "title": "我们提供的价值",
                "lines": ["把“不会用 AI”变成“持续提效”，并把单次试点变成可复制的组织能力"],
            },
        ],
        cols=2,
        body_size=11,
    )
    add_note(
        s,
        6.24,
        "可验证信号：采用率高不等于经营效果高（88% 采用 vs 规模化仍早期），这是“落地方案商”的窗口期。",
    )
    add_footer(s, "下一页：完整提效闭环方案。")

    # B3
    s = base_slide(prs)
    add_header(
        s,
        "CLOSED LOOP",
        "我们的完整闭环方案：非技术团队也能执行",
        "目标不是“做一个聪明 Agent”，而是“建立一条可重复产出价值的业务流水线”。",
        "SLIDE 04 / 10",
    )
    add_cards_block(
        s,
        1.9,
        2.28,
        [
            {
                "title": "1-2. 找问题 + 建模板",
                "lines": ["先锁定最值钱的业务场景", "输出标准输入/输出模板", "降低团队上手门槛"],
            },
            {
                "title": "3-4. 跑执行 + 做复盘",
                "lines": ["通过标准包快速试跑", "统一验收与复盘口径", "确保结果可被管理层理解"],
            },
            {
                "title": "5-6. 沉淀 SOP + 组织扩散",
                "lines": ["把经验固化成 SOP 和资产库", "从个人提效升级为团队提效", "复制到更多业务部门"],
            },
        ],
        cols=3,
        body_size=11,
    )
    add_cards_block(
        s,
        4.32,
        1.35,
        [
            {
                "title": "企业路径",
                "lines": ["199/4999/8999 先验证，再走 29,999 立项，最后进入企业定制与二开交付"],
            },
            {
                "title": "个人路径",
                "lines": ["从 1 个高频任务开始，用教程 + Token + 模板跑通，再复制到更多任务与小团队"],
            },
        ],
        cols=2,
        body_size=10,
    )
    add_metric_row(
        s,
        5.78,
        [
            ("30 天", "拿到首个可汇报样板，降低内部争议"),
            ("90 天", "完成跨岗位复制，形成稳定 SOP"),
            ("180 天", "沉淀行业模板，开始规模化交付"),
        ],
    )
    add_note(s, 6.92, "闭环的本质：让“好结果”可复制，而不是依赖个别高手。", h=0.4)
    add_footer(s, "下一页：竞争优势与护城河。")

    # B4
    s = base_slide(prs)
    add_header(
        s,
        "MOAT",
        "我们的优势：更像“企业 AI 交付系统”，而不只是工具",
        "投资价值来自可复制和可扩张，而不是某一次的演示效果。",
        "SLIDE 05 / 10",
    )
    add_cards_block(
        s,
        1.9,
        2.3,
        [
            {
                "title": "优势 1：落地门槛低",
                "lines": ["开源入口，客户试用成本低", "标准化包先筛选真实需求", "减少盲目大投入"],
            },
            {
                "title": "优势 2：企业可控",
                "lines": ["本地优先和私有化路径", "更适合安全合规要求", "便于组织治理和审计"],
            },
            {
                "title": "优势 3：可规模复制",
                "lines": ["渠道与地区服务商模式放大成交", "模板、工单、SLA 降低交付波动", "经验持续沉淀成资产"],
            },
        ],
        cols=3,
        body_size=11,
    )
    add_cards_block(
        s,
        4.38,
        1.7,
        [
            {
                "title": "为什么不是短期红利",
                "lines": ["核心是“交付系统化”，不是“营销爆点”；系统越跑越多，优势随模板和数据累积增强"],
            },
            {
                "title": "为什么投资可放大",
                "lines": ["资金投入模板资产、渠道能力、交付中台后，新增客户的边际交付效率会持续提升"],
            },
        ],
        cols=2,
        body_size=10,
    )
    add_note(s, 6.2, "替代风险判断：单一 AI 工具容易被替代；“方法+模板+渠道+交付体系”组合更难被替代，且越跑越强。")
    add_footer(s, "下一页：加入对比实验室的精选证据。")

    # New: Selected benchmark evidence from comparison lab
    s = base_slide(prs)
    add_header(
        s,
        "BENCHMARK SELECTED",
        "对比实验室精选：AGIME 在执行、隐私、治理三轴领先",
        "基于 11 款用户向桌面 Agent、4 条评估轴线、25+ 官方资料的可追溯对比。",
        "SLIDE 06 / 10",
    )
    add_metric_row(
        s,
        1.86,
        [
            ("11", "用户向桌面产品（同口径对比）"),
            ("4", "关键评估轴线（执行/隐私/治理/成本）"),
            ("25+", "官方资料交叉验证"),
            ("100%", "证据可追溯（S/D 双编号）"),
        ],
    )
    add_cards_block(
        s,
        3.05,
        2.95,
        [
            {
                "title": "横向矩阵结论（精选）",
                "lines": [
                    "AGIME 在本地离线、文件/系统执行、自动化复用均为“强”",
                    "团队治理维度具备 Team Server、RBAC、审计能力",
                    "支持本地/LAN/私有云部署，兼顾个人与组织场景",
                ],
            },
            {
                "title": "竞品差异（只放关键对照）",
                "lines": [
                    "LobsterAI：自动化强，但公开团队治理信息相对有限",
                    "ChatGPT/Claude：云端研究协作强，但本地控制较弱",
                    "Jan：本地隐私强，但组织治理与协作沉淀能力较弱",
                ],
            },
            {
                "title": "纵向评分口径（可复核）",
                "lines": [
                    "执行深度 = 文件/系统执行 + 自动化调度 + 多步骤闭环",
                    "隐私部署 = 本地离线 + 私有边界 + 模型环境可控",
                    "团队治理 = RBAC/审计 + 协作沉淀 + 扩展生态",
                ],
            },
            {
                "title": "AGIME 三轴评分（精选）",
                "lines": [
                    "执行深度：9.5",
                    "隐私与部署控制：9.6",
                    "团队治理与可扩展：9.1",
                ],
            },
        ],
        cols=2,
        body_size=10,
    )
    add_note(s, 6.12, "注：评分是公开能力映射的相对评估，不替代 PoC 实测；价值在于“先筛选，再深测”。")
    add_footer(s, "下一页：收入模型与增长引擎。")

    # B5
    s = base_slide(prs)
    add_header(
        s,
        "REVENUE",
        "收入模型：3 条引擎共同驱动增长",
        "先用标准包获得现金流和高意向客户，再用渠道放大，再用定制提高利润。",
        "SLIDE 07 / 10",
    )
    add_cards_block(
        s,
        1.9,
        2.18,
        [
            {
                "title": "引擎 A：标准服务包",
                "lines": ["199 / 4999 / 8999 一次性交付", "快速回款，筛选真实客户", "构建立项前证据链"],
            },
            {
                "title": "引擎 B：渠道与生态收入",
                "lines": ["代理/区域/生态准入费", "渠道分成与平台协同费", "保证金形成现金流安全垫"],
            },
            {
                "title": "引擎 C：企业定制交付",
                "lines": ["29,999 立项包锁定边界", "MCP Skills 和二开项目提客单", "高毛利来自标准化交付能力"],
            },
        ],
        cols=3,
        body_size=10,
    )
    add_cards_block(
        s,
        4.24,
        1.35,
        [
            {
                "title": "赚钱节奏",
                "lines": ["标准包解决现金流，渠道收入解决规模化，定制交付解决利润率，三者形成正循环"],
            },
            {
                "title": "核心投资价值",
                "lines": ["同一套方法可在不同区域、行业复制，收入结构不依赖单一大客户"],
            },
        ],
        cols=2,
        body_size=10,
    )
    add_cards_block(
        s,
        5.74,
        1.36,
        [
            {
                "title": "回款速度",
                "lines": ["一次性产品包回款更快", "减少长期陪跑现金流压力", "先筛选再重投，控制坏项目比例"],
            },
            {
                "title": "基准收入锚点",
                "lines": ["当前服务基准月收入：¥381,910", "叠加渠道确认与分成收入", "形成“现金流+利润”双引擎"],
            },
            {
                "title": "上行空间",
                "lines": ["标准包扩大成交面", "立项包提升定制签约概率", "渠道复制提升区域扩张速度"],
            },
        ],
        cols=3,
        body_size=9,
    )
    add_footer(s, "下一页：融资后 12 个月执行计划。")

    # B6
    s = base_slide(prs)
    add_header(
        s,
        "PLAN",
        "融资用途与 12 个月目标：把“方法”做成“规模”",
        "资金全部投向可复用资产，确保每一笔投入都提升后续交付效率和转化效率。",
        "SLIDE 08 / 10",
    )
    add_cards_block(
        s,
        1.9,
        2.3,
        [
            {
                "title": "资金怎么用",
                "lines": [
                    "40%：产品与交付中台（模板库、工单、指标系统）",
                    "30%：渠道体系（代理/地区服务商拓展与管理）",
                    "20%：行业解决方案（MCP Skills 与场景包）",
                    "10%：合规与风控能力建设",
                ],
            },
            {
                "title": "12 个月里程碑",
                "lines": [
                    "形成标准化行业方案包并可复制交付",
                    "建立稳定渠道网络并跑通区域协同",
                    "提升标准包->立项、立项->定制转化率",
                    "实现收入结构更均衡、现金流更稳健",
                ],
            },
        ],
        cols=2,
        body_size=10,
    )
    add_metric_row(
        s,
        4.36,
        [
            ("目标 1", "把提效闭环在更多行业复制"),
            ("目标 2", "把收入从“项目型”升级到“体系型”"),
            ("目标 3", "把交付能力沉淀为长期护城河"),
        ],
    )
    add_cards_block(
        s,
        5.55,
        1.35,
        [
            {
                "title": "收益逻辑",
                "lines": [
                    "早期看现金流安全：标准包与渠道收入托底",
                    "中期看利润抬升：定制与行业方案放大客单",
                    "长期看估值提升：交付体系资产化与可复制网络",
                ],
            },
            {
                "title": "执行纪律",
                "lines": [
                    "每月复盘：成交结构、回款周期、客诉率",
                    "每季度复盘：渠道质量、区域冲突、项目毛利",
                    "每半年复盘：行业模板复用率和转化效率",
                ],
            },
        ],
        cols=2,
        body_size=9,
    )
    add_footer(s, "结论：在 AI 逻辑焦虑市场中，我们提供“看得懂、做得到、可复制、可放大”的标准答案。")

    # B7
    s = base_slide(prs)
    add_header(
        s,
        "INVESTOR FAQ",
        "常见质疑与回应",
        "聚焦三类核心质疑：可执行性、盈利性、可持续性。",
        "SLIDE 09 / 10",
    )
    add_cards_block(
        s,
        1.9,
        2.32,
        [
            {
                "title": "质疑 1：这是不是短期风口？",
                "lines": ["需求不是“买模型”，而是“做结果”", "我们解决的是长期经营问题：效率与复用", "组织级流程改造周期更长、黏性更高"],
            },
            {
                "title": "质疑 2：会不会只靠老板关系单？",
                "lines": ["用标准包做可规模获客", "渠道与地区服务商帮助持续扩张", "收入结构不是“单大单驱动”"],
            },
            {
                "title": "质疑 3：为什么别人不能很快复制？",
                "lines": ["单工具易复制，交付系统难复制", "模板、SOP、工单、数据资产持续累积", "越交付越强，形成复利型护城河"],
            },
        ],
        cols=3,
        body_size=10,
    )
    add_cards_block(
        s,
        4.38,
        1.78,
        [
            {
                "title": "投资判断标准",
                "lines": ["是否能持续提升客户效率结果", "是否能在更多区域和行业复制成交", "是否能让收入结构越来越健康"],
            },
            {
                "title": "投后跟踪动作",
                "lines": ["按季度看转化率和回款，而非单月波动", "跟踪模板复用率和项目交付稳定性", "优先支持“可复制能力”建设投入"],
            },
        ],
        cols=2,
        body_size=10,
    )
    add_note(s, 6.26, "结论：AGIME 的说服力不在“说得多好”，而在“能否持续交付并复制增长”。")
    add_footer(s, "附：完整来源、证据和审计。")

    # S8
    s = base_slide(prs)
    add_header(
        s,
        "SOURCES & AUDIT",
        "资料来源与 ID 审计",
        "规则：Txx=证据结论；Rxx=原始链接。以下内容同步 financing-plan.html 当前版本。",
        "SLIDE 10 / 10",
    )

    left_lines = [
        "T01 Gartner：2026 AI支出 2.52万亿美元（YoY +44%）",
        "T02 McKinsey：88%组织已使用AI",
        "T03 McKinsey：62%试验agents、39%有EBIT影响",
        "T04 OECD：2025企业AI采用率20.2%",
        "T05 OECD：2025超三分之一人使用GenAI",
        "T06 OpenAI：100万+ business customers，700万+ seats",
        "T07 WEF：到2030年39%核心技能将变化",
        "T08 Stack Overflow：80%使用AI，信任29%，66%返工",
        "T09 AGIME评分：执行9.5、隐私9.6、治理9.1",
        "T10 AGIME能力：本地优先、团队协作、企业部署",
        "T11 商业化定价：199/4999/8999/29,999+定制",
        "T12 Gartner：2025 GenAI支出 6440亿美元（+76.4%）",
        "T13 基准月收入381,910（60/20/10/3）",
    ]
    right_lines = [
        "R01 https://www.gartner.com/.../2026-1-15-...-2-point-5-trillion-...",
        "R02 https://www.mckinsey.com/.../the-state-of-ai",
        "R03 https://www.oecd.org/digital/artificial-intelligence/",
        "R04 https://openai.com/business/.../state-of-enterprise-ai-2025-report/",
        "R05 https://www.weforum.org/.../future-of-jobs-report-2025/...",
        "R06 https://stackoverflow.blog/2025/12/29/...developer-survey...",
        "R07 product-comparison.html#vertical",
        "R08 https://github.com/jsjm1986/AGIME/blob/main/README.md",
        "R09 business-channel.html#pricing",
        "R10 https://www.gartner.com/.../2025-03-31-...-644-billion-in-2025",
        "R11 business-channel.html#calculator",
        "审计说明：预测与测算为经营假设，不构成收益承诺。",
    ]
    add_cards_block(
        s,
        1.9,
        4.95,
        [
            {"title": "证据结论（T01-T13）", "lines": left_lines},
            {"title": "来源映射（R01-R11）", "lines": right_lines},
        ],
        cols=2,
        body_size=8,
    )
    add_note(s, 6.95, "复核时间：2026-02-24（UTC+8）｜证据链与链接链一一对应，可供尽调抽查。", h=0.4)


def build_basic_slides_v3(prs: Presentation) -> None:
    # 01 Cover thesis
    s = base_slide(prs)
    add_header(
        s,
        "NON-TECH DECISION · 2026-02-24",
        "AGIME：把 AI 焦虑变成可复制收益",
        "面向非技术投资决策：我们不是卖模型概念，而是交付可复用的效率结果。",
        "SLIDE 01 / 16",
    )
    add_metric_row(
        s,
        1.86,
        [
            ("$2.52T", "2026 全球 AI 支出预测"),
            ("88%", "组织已在业务中使用 AI"),
            ("3 条", "服务包 + 渠道 + 定制收入引擎"),
            ("6 步", "诊断->模板->执行->复盘->SOP->扩散"),
        ],
    )
    add_cards_block(
        s,
        3.02,
        2.7,
        [
            {
                "title": "投资人最关心的 3 个问题",
                "lines": [
                    "这是不是短期风口，还是长期需求？",
                    "收入是否可复制，不依赖单一大客户？",
                    "团队优势是否会被快速替代？",
                ],
            },
            {
                "title": "我们的直给答案",
                "lines": [
                    "客户买的是“稳定提效结果”，不是模型噱头",
                    "收入来自组合结构，天然抗单点波动",
                    "壁垒在交付系统化，越跑越强",
                ],
            },
        ],
        cols=2,
        body_size=11,
    )
    add_note(s, 5.84, "一句话：AGIME 在 AI 逻辑焦虑市场中，给出“看得懂、做得到、可复制”的落地答案。")
    add_footer(s, "下一页：官网口号与产品名的商业含义。")

    # 02 Homepage narrative and product naming
    s = base_slide(prs)
    add_header(
        s,
        "HOMEPAGE NARRATIVE",
        "产品名与口号说明：AGIME 是协作系统，不是聊天工具",
        "官网首屏原文聚焦 Human × AI × Team × Ecosystem，强调从个人提效到组织复用。",
        "SLIDE 02 / 16",
    )
    add_cards_block(
        s,
        1.9,
        3.6,
        [
            {
                "title": "首页口号（原文）",
                "lines": [
                    "不只是一个 AI 工具，而是一套协作新范式",
                    "从“人与 AI”升级到“团队与 AI”再到“生态协同”",
                    "把个人能力、团队能力、外部协作放进同一执行系统",
                ],
            },
            {
                "title": "AGIME 名称叙事",
                "lines": [
                    "AGIME = AI + Me 的组织化延展",
                    "Human 负责决策，AI 负责执行",
                    "Team 负责治理复用，Ecosystem 负责跨组织交付",
                ],
            },
            {
                "title": "协作链路（原文）",
                "lines": [
                    "个人沉淀能力 -> 团队共享复用",
                    "Agent 云端执行 -> 对外协作交付",
                    "把“个人会做”升级为“团队都会做”",
                ],
            },
            {
                "title": "核心价值（原文）",
                "lines": [
                    "永久免费开源，降低试错门槛",
                    "本地优先与数据可控，适配企业安全要求",
                    "支持 Team Server 与私有化，便于规模部署",
                ],
            },
        ],
        cols=2,
        body_size=10,
    )
    add_metric_row(
        s,
        5.68,
        [
            ("个人层", "文件处理、自动化、流程沉淀"),
            ("团队层", "共享 Skills/MCP、权限、审计"),
            ("生态层", "Portal 对外协作与统一交付"),
        ],
    )
    add_footer(s, "下一页：为什么现在是“落地方案年”。")

    # 03 AI trend and timing
    s = base_slide(prs)
    add_header(
        s,
        "AI TREND",
        "趋势窗口：预算高增长，但企业仍缺“稳定落地能力”",
        "市场已从工具采用转向结果治理，买方关注“能否长期稳定提效”。",
        "SLIDE 03 / 16",
    )
    add_metric_row(
        s,
        1.86,
        [
            ("$2.52T", "2026 全球 AI 支出"),
            ("$644B", "2025 GenAI 支出（+76.4%）"),
            ("62%", "企业在试验 Agents"),
            ("39%", "有 EBIT 影响（规模化仍早期）"),
        ],
    )
    add_cards_block(
        s,
        3.04,
        2.75,
        [
            {
                "title": "趋势变化",
                "lines": [
                    "工作流时代解决“能不能用”",
                    "Agent 时代解决“能不能做”",
                    "当前进入“能不能规模化提效”阶段",
                ],
            },
            {
                "title": "买方真实诉求",
                "lines": [
                    "不是更炫功能，而是更稳结果",
                    "需要可审计、可治理、可复制",
                    "愿为确定性交付付费",
                ],
            },
            {
                "title": "窗口判断",
                "lines": [
                    "模型能力趋同，交付能力重估",
                    "先做确定性，再做规模化",
                    "交付系统型公司进入估值窗口",
                ],
            },
        ],
        cols=3,
        body_size=10,
    )
    add_note(s, 5.9, "结论：机会不在“再造一个模型入口”，而在“把执行与治理打包交付”。")
    add_footer(s, "下一页：市场需求与痛点缺口。")

    # 04 Market gap and demand
    s = base_slide(prs)
    add_header(
        s,
        "MARKET GAP",
        "市场需求：大家都在买 AI，但多数团队不会系统提效",
        "高采用率不等于高经营效果，企业缺的是“从试点到规模”的中间层。",
        "SLIDE 04 / 16",
    )
    add_cards_block(
        s,
        1.9,
        2.5,
        [
            {
                "title": "现实痛点",
                "lines": [
                    "会用 AI，但结果波动大",
                    "试点多，组织复用少",
                    "投入后难复盘 ROI",
                ],
            },
            {
                "title": "信任痛点",
                "lines": [
                    "80% 在用 AI，但信任仅 29%",
                    "66% 花时间修复“差一点正确”结果",
                    "缺统一验收与质量门槛",
                ],
            },
            {
                "title": "市场错位",
                "lines": [
                    "供给偏工具，需求偏结果",
                    "客户缺落地方法，不缺按钮",
                    "安全治理顾虑拉长采购周期",
                ],
            },
        ],
        cols=3,
        body_size=10,
    )
    add_cards_block(
        s,
        4.56,
        1.55,
        [
            {
                "title": "目标客户画像",
                "lines": ["正在探索 AI 转型，但“方向不清、团队不会、成本不可控”的企业与团队"],
            },
            {
                "title": "我们的切入价值",
                "lines": ["把“不会用 AI”变成“持续提效”，再把单点样板复制成组织能力"],
            },
        ],
        cols=2,
        body_size=10,
    )
    add_note(s, 6.22, "机会判断：谁能稳定交付结果，谁就能持续拿预算。")
    add_footer(s, "下一页：我们如何把提效做成闭环。")

    # 05 Closed-loop solution
    s = base_slide(prs)
    add_header(
        s,
        "CLOSED LOOP",
        "完整闭环方案：非技术团队也能执行",
        "目标不是“演示聪明”，而是“建立可重复产出价值的业务流水线”。",
        "SLIDE 05 / 16",
    )
    add_cards_block(
        s,
        1.9,
        2.32,
        [
            {
                "title": "1-2. 找问题 + 建模板",
                "lines": ["先锁定最值钱场景", "标准化输入/输出模板", "降低上手门槛"],
            },
            {
                "title": "3-4. 跑执行 + 做复盘",
                "lines": ["标准包先试跑", "统一验收与复盘口径", "结果可被管理层理解"],
            },
            {
                "title": "5-6. 沉淀 SOP + 扩散",
                "lines": ["经验固化为 SOP 与资产库", "从个人提效升级为团队提效", "再复制到更多业务部门"],
            },
        ],
        cols=3,
        body_size=10,
    )
    add_cards_block(
        s,
        4.38,
        1.34,
        [
            {
                "title": "企业路径",
                "lines": ["199/4999/8999 先验证，再走 29,999 立项，最后进入企业定制与二开交付"],
            },
            {
                "title": "个人路径",
                "lines": ["从 1 个高频任务起步：教程 + Token + 模板跑通，再复制到更多任务与小团队"],
            },
        ],
        cols=2,
        body_size=9,
    )
    add_metric_row(
        s,
        5.8,
        [
            ("30 天", "拿到首个可汇报样板"),
            ("90 天", "形成跨岗位稳定 SOP"),
            ("180 天", "沉淀行业模板并规模化"),
        ],
    )
    add_footer(s, "下一页：Team Server 如何把个人能力变组织能力。")

    # 06 Team Server
    s = base_slide(prs)
    add_header(
        s,
        "TEAM SERVER",
        "组织化关键：Team Server 把“人效”升级为“组织效”",
        "团队不是每个人各配一套 AI，而是共享一个可治理、可审计、可持续执行的中枢。",
        "SLIDE 06 / 16",
    )
    add_cards_block(
        s,
        1.9,
        3.9,
        [
            {
                "title": "云端 Team Agent 执行中枢",
                "lines": [
                    "任务放到云端持续执行，不占个人电脑",
                    "自动调用团队 Skills、MCP 工具和文档知识",
                    "执行结果可直接复用给全团队",
                ],
            },
            {
                "title": "共享与沉淀",
                "lines": [
                    "一人沉淀，全员复用",
                    "新人入组当天可开工",
                    "减少反复问、反复试",
                ],
            },
            {
                "title": "治理与风控",
                "lines": [
                    "角色、分组、权限、审核可配置",
                    "关键动作可追踪、可审计",
                    "适配敏感流程与合规要求",
                ],
            },
            {
                "title": "内外协同",
                "lines": [
                    "Portal 对外统一发布资料与答疑",
                    "跨地团队统一版本与流程",
                    "支持 LAN 与云端协同部署",
                ],
            },
        ],
        cols=2,
        body_size=10,
    )
    add_note(s, 5.98, "关系说明：Team 不是替代个人版，而是在个人能力之上增加“共享、治理、云端执行”。")
    add_metric_row(
        s,
        6.62,
        [
            ("个人层", "本地执行与流程沉淀"),
            ("团队层", "共享复用与治理审计"),
            ("生态层", "对外协作交付窗口"),
        ],
    )

    # 07 Scenarios
    s = base_slide(prs)
    add_header(
        s,
        "SCENARIO VALUE",
        "价值落地：不是讲故事，而是覆盖高频可计量场景",
        "从文档处理到自动化执行，再到团队协作交付，形成可持续人效提升。",
        "SLIDE 07 / 16",
    )
    add_cards_block(
        s,
        1.9,
        4.86,
        [
            {
                "title": "文档与表格批处理",
                "lines": ["批量读取合同/PDF/Excel", "自动提取关键字段与汇总", "减少重复人工核对"],
            },
            {
                "title": "桌面自动化",
                "lines": ["模拟点击、键盘输入、窗口控制", "跨系统重复动作自动执行", "降低低价值手工工时"],
            },
            {
                "title": "定时与持续执行",
                "lines": ["教一次生成可复用规范", "支持定时定量自动运行", "7×24 连续执行不遗漏"],
            },
            {
                "title": "调研与信息整理",
                "lines": ["自动浏览网页采集信息", "按模板整理为结构化输出", "用于竞品、行业与客户研究"],
            },
            {
                "title": "数据分析与报告",
                "lines": ["辅助分析数据与发现异常", "生成图表和摘要报告", "提升管理层决策效率"],
            },
            {
                "title": "团队协作复用",
                "lines": ["把个人流程沉淀为团队模板", "跨部门统一执行与验收标准", "持续积累组织知识资产"],
            },
        ],
        cols=3,
        body_size=9,
    )
    add_footer(s, "下一页：对比实验室样本与方法。")

    # 08 Benchmark scope
    s = base_slide(prs)
    add_header(
        s,
        "BENCHMARK SCOPE",
        "对比实验室：样本、口径、证据三件事先讲清楚",
        "我们只对比“用户向桌面 Agent”，避免把编程 Agent 与办公 Agent 混为一谈。",
        "SLIDE 08 / 16",
    )
    add_metric_row(
        s,
        1.86,
        [
            ("11", "用户向桌面产品"),
            ("4", "关键评估轴线"),
            ("25+", "官方资料交叉验证"),
            ("100%", "证据可追溯"),
        ],
    )
    add_cards_block(
        s,
        3.04,
        3.08,
        [
            {
                "title": "横向矩阵看什么",
                "lines": [
                    "本地离线能力",
                    "文件/系统执行能力",
                    "自动化与工作流能力",
                    "团队治理与私有部署能力",
                ],
            },
            {
                "title": "纵向评分口径",
                "lines": [
                    "执行深度：执行链路闭环能力",
                    "隐私部署：本地与数据边界控制",
                    "团队治理：RBAC/审计/协作沉淀",
                ],
            },
            {
                "title": "核心样本（精选）",
                "lines": [
                    "AGIME / LobsterAI / QoderWork / OpenClaw",
                    "ChatGPT Desktop / Claude Desktop",
                    "Copilot / Perplexity / Jan 等",
                ],
            },
            {
                "title": "审计结构",
                "lines": [
                    "Sxx：结论证据（聚合）",
                    "Dxx：原始文档（逐条）",
                    "支持尽调时逐项复核",
                ],
            },
        ],
        cols=2,
        body_size=10,
    )
    add_footer(s, "下一页：对比结果（横向 + 纵向）精选。")

    # 09 Benchmark result
    s = base_slide(prs)
    add_header(
        s,
        "BENCHMARK RESULT",
        "对比结论：AGIME 在执行、隐私、治理三轴保持领先",
        "分值用于快速筛选，不替代 PoC；价值在于更快确定验证路径。",
        "SLIDE 09 / 16",
    )
    add_metric_row(
        s,
        1.86,
        [
            ("9.5", "执行深度"),
            ("9.6", "隐私与部署控制"),
            ("9.1", "团队治理与可扩展"),
        ],
    )
    add_cards_block(
        s,
        3.05,
        3.18,
        [
            {
                "title": "执行深度对照（精选）",
                "lines": [
                    "AGIME 9.5 / LobsterAI 9.0 / OpenClaw 8.8",
                    "Claude 7.6 / QoderWork 7.5 / ChatGPT 7.4",
                    "结论：AGIME 更接近“可执行系统”",
                ],
            },
            {
                "title": "隐私部署对照（精选）",
                "lines": [
                    "AGIME 9.6 / Jan 9.1 / OpenClaw 8.6",
                    "LobsterAI 8.2 / QoderWork 7.5 / Perplexity 5.6",
                    "结论：本地优先 + 私有路径是企业关键门槛",
                ],
            },
            {
                "title": "治理扩展对照（精选）",
                "lines": [
                    "AGIME 9.1 / Perplexity 8.8 / ChatGPT 8.4",
                    "Claude 8.0 / Copilot 7.8 / QoderWork 7.6",
                    "结论：AGIME 团队治理能力接近企业系统形态",
                ],
            },
        ],
        cols=3,
        body_size=9,
    )
    add_note(s, 6.35, "行动建议：先用矩阵筛选，再对 2-3 个候选做 2 周 PoC，避免盲目长周期试错。")
    add_footer(s, "下一页：为什么这种优势能形成护城河。")

    # 10 Moat
    s = base_slide(prs)
    add_header(
        s,
        "MOAT",
        "护城河：方法 + 模板 + 渠道 + 交付体系的复利效应",
        "单一工具易复制，系统化交付难复制，且会随交付数据持续变强。",
        "SLIDE 10 / 16",
    )
    add_cards_block(
        s,
        1.9,
        2.36,
        [
            {
                "title": "优势 1：低门槛验证",
                "lines": ["开源入口降低试错成本", "标准包先筛选真实需求", "减少无效投入"],
            },
            {
                "title": "优势 2：企业可控",
                "lines": ["本地优先与私有化路径", "适配安全与合规", "便于治理与审计"],
            },
            {
                "title": "优势 3：可规模复制",
                "lines": ["模板化交付降低波动", "渠道网络放大成交", "经验沉淀为长期资产"],
            },
        ],
        cols=3,
        body_size=10,
    )
    add_cards_block(
        s,
        4.44,
        1.66,
        [
            {
                "title": "为什么不是短期红利",
                "lines": ["核心价值是“交付系统化”而非营销爆点；系统越跑越多，优势随模板和数据持续增强。"],
            },
            {
                "title": "为什么可投资放大",
                "lines": ["投入模板资产、渠道能力、交付中台后，新增客户边际交付成本下降，利润率和确定性同步提升。"],
            },
        ],
        cols=2,
        body_size=9,
    )
    add_note(s, 6.2, "替代风险判断：产品可被模仿，但“结果可复制”的交付体系很难被短期复制。")
    add_footer(s, "下一页：收入模型总览。")

    # 11 Revenue model
    s = base_slide(prs)
    add_header(
        s,
        "REVENUE ENGINE",
        "收入模型：3 条引擎驱动“现金流 + 规模 + 利润”",
        "先做标准化回款，再做渠道放大，再做高客单定制。",
        "SLIDE 11 / 16",
    )
    add_cards_block(
        s,
        1.9,
        2.24,
        [
            {
                "title": "引擎 A：标准服务包",
                "lines": ["199 / 4999 / 8999 一次性交付", "快速回款并筛选高意向", "构建立项前证据链"],
            },
            {
                "title": "引擎 B：渠道与生态收入",
                "lines": ["代理/区域/生态准入费", "标准包返佣与平台协同费", "保证金增强履约与价格纪律"],
            },
            {
                "title": "引擎 C：企业定制交付",
                "lines": ["29,999 立项包锁边界", "MCP Skills 与 AGIME 二开提客单", "高毛利来自标准化交付能力"],
            },
        ],
        cols=3,
        body_size=10,
    )
    add_cards_block(
        s,
        4.28,
        1.38,
        [
            {
                "title": "赚钱节奏",
                "lines": ["标准包保现金流，渠道保规模，定制保利润率；三者形成正循环。"],
            },
            {
                "title": "投资价值",
                "lines": ["同一方法可跨区域、跨行业复制，收入结构不依赖单一大客户。"],
            },
        ],
        cols=2,
        body_size=10,
    )
    add_metric_row(
        s,
        5.84,
        [
            ("¥381,910", "当前服务基准月收入"),
            ("80+ / 20+ / 10+", "199 / 4999 / 8999 建议月成交"),
            ("29,999+", "立项后转入定制客单层"),
        ],
    )
    add_footer(s, "下一页：服务产品化定价（含教程 + Token 组合）。")

    # 12 Pricing details
    s = base_slide(prs)
    add_header(
        s,
        "PRICING DETAIL",
        "服务产品化与定价：标准化输出，低人力占比",
        "定价策略：先用一次性包筛选有效需求，再转入 29,999 立项与企业定制。",
        "SLIDE 12 / 16",
    )
    add_cards_block(
        s,
        1.9,
        4.82,
        [
            {
                "title": "199 诊断启动包｜¥199",
                "lines": [
                    "标准教程库 + $99 Token 启动包 + 评分问卷 + 任务模板",
                    "系统自动生成诊断报告，适合线索筛选",
                    "低人工参与，标准化交付",
                ],
            },
            {
                "title": "4999 部署验证包｜¥4,999",
                "lines": [
                    "私有化部署脚本与指南 + 去 Logo 商用授权 + $299 Token",
                    "标准配置模板 + 验收文档",
                    "一次性验证包，不含定制开发",
                ],
            },
            {
                "title": "8999 业务试点包｜¥8,999",
                "lines": [
                    "试点模板库 + 团队自助教程 + $499 Token + 自动复盘报表",
                    "输出 ROI 试点评估与需求模板",
                    "一次性交付，不含定制开发",
                ],
            },
            {
                "title": "29,999 定制立项包｜¥29,999",
                "lines": [
                    "需求采集模板 + 技术架构蓝图 + $999 Token + 里程碑报价模板",
                    "可 100% 抵扣后续定制合同首款",
                    "立项后按里程碑进入 MCP/二开开发",
                ],
            },
        ],
        cols=2,
        body_size=9,
    )
    add_note(s, 6.82, "授权边界：开源代码仍遵循 Apache-2.0，收费项对应部署、去标授权、培训与定制开发交付。")
    add_footer(s, "下一页：渠道生态合作的费用与收益机制。")

    # 13 Partner economics
    s = base_slide(prs)
    add_header(
        s,
        "CHANNEL & ECOSYSTEM",
        "合作生态：费用、分成、支持三件事透明化",
        "伙伴既是客户也是增长引擎，合作目标是“会做、能卖、可持续赚钱”。",
        "SLIDE 13 / 16",
    )
    add_cards_block(
        s,
        1.9,
        2.7,
        [
            {
                "title": "代理销售伙伴",
                "lines": [
                    "准入费 ¥59,800；保证金 ¥50,000；最低采购 ¥199,999/年",
                    "返佣：199包30% / 4999包40% / 8999包45%",
                    "定制销售佣金：18%-25%",
                ],
            },
            {
                "title": "地区服务商（省/市/区）",
                "lines": [
                    "准入费：区 ¥89,800 / 市 ¥199,800 / 省 ¥398,000 起",
                    "交付分成：35%-50%；总部收 8% 平台协同费",
                    "省级保证金 ¥200,000（市100,000 / 区50,000）",
                ],
            },
            {
                "title": "技术生态伙伴",
                "lines": [
                    "准入费 ¥9,800；保证金 ¥5,000；最低采购 ¥19,999/年",
                    "市场分成：80%（伙伴）/20%（平台）",
                    "提供 SDK/API 文档、示例代码与联合发布渠道",
                ],
            },
        ],
        cols=3,
        body_size=9,
    )
    add_cards_block(
        s,
        4.78,
        1.68,
        [
            {
                "title": "合作流程（4 步）",
                "lines": ["选择身份 -> 签约认证 -> 标准包成交 -> 转入立项/定制；结算按回款与里程碑执行。"],
            },
            {
                "title": "总部支持边界",
                "lines": ["标准方案库、报价模板、工单系统、培训认证、SLA 支持；重点是“标准化支持优先”。"],
            },
            {
                "title": "可读收益模型",
                "lines": ["伙伴收益 = 销售返佣 + 交付分成 + 生态分成 + 阶梯激励，路径明确且可测算。"],
            },
        ],
        cols=3,
        body_size=9,
    )
    add_footer(s, "下一页：财务模型与关键经营假设。")

    # 14 Financial model
    s = base_slide(prs)
    add_header(
        s,
        "FINANCIAL MODEL",
        "财务模型：先做确定性，再做规模化",
        "收入按“服务 + 渠道 + 定制”三池拆分，便于按月复盘转化效率与交付利润。",
        "SLIDE 14 / 16",
    )
    add_metric_row(
        s,
        1.86,
        [
            ("¥381,910", "服务池基准月收入（当前模型）"),
            ("60/20/10/3", "199/4999/8999/定制 月基准订单"),
            ("3 收入池", "服务直签 + 渠道确认 + 渠道分成"),
            ("T+15", "标准包回款后结算参考周期"),
        ],
    )
    add_cards_block(
        s,
        3.05,
        3.2,
        [
            {
                "title": "服务收入池（直签）",
                "lines": [
                    "199/4999/8999 一次性标准包提供现金流底座",
                    "29,999 立项包锁定边界，提高定制成功率",
                    "核心指标：订单量、交付工时、毛利率、客诉率",
                ],
            },
            {
                "title": "渠道收入池（代理/区域/生态）",
                "lines": [
                    "准入费与最低采购额形成基础现金流",
                    "标准包返佣与平台协同费形成可持续分配机制",
                    "核心指标：达标率、冲突率、回款质量",
                ],
            },
            {
                "title": "融资后投入比例（12 个月）",
                "lines": [
                    "40% 产品与交付中台（模板库/工单/指标）",
                    "30% 渠道体系与区域协同",
                    "20% 行业方案（MCP Skills 场景包）+ 10% 合规风控",
                ],
            },
            {
                "title": "管理原则",
                "lines": [
                    "按月复盘成交结构与回款周期",
                    "按季度复盘渠道质量与项目毛利",
                    "避免“收入增长但人效下降”的假繁荣",
                ],
            },
        ],
        cols=2,
        body_size=9,
    )
    add_footer(s, "下一页：融资用途与 12 个月执行路线。")

    # 15 Execution plan
    s = base_slide(prs)
    add_header(
        s,
        "EXECUTION PLAN",
        "融资后 12 个月计划：把方法沉淀成规模能力",
        "执行核心：每一笔投入都服务于“可复制成交 + 可复制交付”。",
        "SLIDE 15 / 16",
    )
    add_cards_block(
        s,
        1.9,
        2.48,
        [
            {
                "title": "0-3 个月：标准化加固",
                "lines": [
                    "完善 199/4999/8999 交付清单与验收模板",
                    "上线统一报价、工单、复盘机制",
                    "确保“低人工占比”的交付纪律",
                ],
            },
            {
                "title": "4-8 个月：渠道规模化",
                "lines": [
                    "重点拓展代理与地区服务商网络",
                    "按省/市/区推进区域协同与冲突治理",
                    "形成可复制的伙伴 onboarding 机制",
                ],
            },
            {
                "title": "9-12 个月：行业深化",
                "lines": [
                    "沉淀 2-3 个行业 MCP Skills 方案包",
                    "提升立项包->定制签约率与客单价",
                    "形成“模板资产 + 区域网络”双轮驱动",
                ],
            },
        ],
        cols=3,
        body_size=9,
    )
    add_cards_block(
        s,
        4.56,
        1.9,
        [
            {
                "title": "管理看板（投后重点）",
                "lines": [
                    "成交结构：标准包/立项包/定制占比",
                    "财务质量：回款周期、坏账、毛利率",
                    "交付质量：客诉率、复用率、SLA 达成率",
                ],
            },
            {
                "title": "阶段目标（示例）",
                "lines": [
                    "月度 199 包 80+，4999 包 20+，8999 包 10+",
                    "季度定制立项签约 3+",
                    "收入结构持续去单点化",
                ],
            },
        ],
        cols=2,
        body_size=9,
    )
    add_note(s, 6.54, "投后目标不是“追热点”，而是持续提升“可复制成交 + 可复制交付”的系统能力。")
    add_footer(s, "下一页：常见质疑、风险对冲与来源审计。")

    # 16 FAQ + Sources
    s = base_slide(prs)
    add_header(
        s,
        "FAQ · SOURCES",
        "常见质疑、风险对冲与来源审计",
        "本页可作为尽调入口：先看质疑回应，再按证据 ID 抽查原始来源。",
        "SLIDE 16 / 16",
    )
    add_cards_block(
        s,
        1.9,
        3.9,
        [
            {
                "title": "常见质疑与回应（精选）",
                "lines": [
                    "Q：是不是短期风口？A：我们解决的是长期经营效率问题",
                    "Q：会不会依赖关系单？A：标准包 + 渠道体系降低单点依赖",
                    "Q：会被快速复制吗？A：工具可复制，交付系统难复制",
                    "Q：数据安全吗？A：本地优先 + 私有化路径 + 权限审计",
                ],
            },
            {
                "title": "主要风险与对冲",
                "lines": [
                    "ROI 波动：先标准包验证再立项重投",
                    "交付返工：统一验收与关键节点复核",
                    "渠道失控：统一模板、工单、冲突裁定机制",
                    "组织依赖：SOP 与模板库降低关键人风险",
                ],
            },
            {
                "title": "来源审计（外部）",
                "lines": [
                    "Gartner：2026 AI 支出、2025 GenAI 支出",
                    "McKinsey：AI 采用与 Agent 规模化阶段",
                    "WEF / Stack Overflow：技能变化与信任/返工数据",
                    "OpenAI/OECD：企业使用规模与趋势信号",
                ],
            },
            {
                "title": "来源审计（内部）",
                "lines": [
                    "官网：产品口号、Team Server 叙事与场景",
                    "产品对比页：横向矩阵、纵向评分、S/D 证据结构",
                    "商业化页：定价包、渠道费用、分成与结算规则",
                    "注：预测与测算为经营假设，不构成收益承诺",
                ],
            },
        ],
        cols=2,
        body_size=9,
    )
    add_note(s, 5.96, "尽调建议：按“结论->证据->原始链接”三步复核，优先验证转化率、回款质量、交付复用率。")
    add_metric_row(
        s,
        6.55,
        [
            ("16 页", "非技术决策完整叙事版"),
            ("3 页面", "官网 + 对比实验 + 商业化同步"),
            ("可执行", "可直接用于投资人沟通"),
        ],
    )


def main() -> None:
    prs = new_presentation()
    build_basic_slides_v3(prs)

    out_dir = Path("wwwweb")
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / "AGIME_financing_non_tech_v3_16slides_2026-02-24.pptx"
    prs.save(str(out_path))
    print(f"Generated: {out_path}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
