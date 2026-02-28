from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BG=RGBColor(0x0D,0x1B,0x2A); DGREEN=RGBColor(0x1A,0x3C,0x34); GOLD=RGBColor(0xC9,0xA8,0x4C)
GREEN=RGBColor(0x2E,0xCC,0x71); WHITE=RGBColor(0xFF,0xFF,0xFF); LGRAY=RGBColor(0xCC,0xCC,0xCC)
RED=RGBColor(0xE7,0x4C,0x3C); CARD=RGBColor(0x16,0x2A,0x3A); MCARD=RGBColor(0x0F,0x24,0x33)
ORANGE=RGBColor(0xF3,0x9C,0x12)

W,H=Cm(25.4),Cm(14.3)
prs=Presentation(); prs.slide_width=W; prs.slide_height=H

def S(): return prs.slides.add_slide(prs.slide_layouts[6])

def bg(sl,c=BG):
    r=sl.shapes.add_shape(1,0,0,W,H); r.fill.solid(); r.fill.fore_color.rgb=c; r.line.fill.background()

def R(sl,x,y,w,h,c):
    r=sl.shapes.add_shape(1,Cm(x),Cm(y),Cm(w),Cm(h)); r.fill.solid(); r.fill.fore_color.rgb=c; r.line.fill.background(); return r

def T(sl,text,x,y,w,h,sz=14,c=WHITE,bold=False,align=PP_ALIGN.LEFT,italic=False):
    tb=sl.shapes.add_textbox(Cm(x),Cm(y),Cm(w),Cm(h)); tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=align; run=p.add_run(); run.text=text
    run.font.size=Pt(sz); run.font.color.rgb=c; run.font.bold=bold; run.font.italic=italic; return tb

def TL(sl,lines,x,y,w,h):
    """lines: list of (text,sz,color,bold)"""
    tb=sl.shapes.add_textbox(Cm(x),Cm(y),Cm(w),Cm(h)); tf=tb.text_frame; tf.word_wrap=True
    for i,(text,sz,c,bold) in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        run=p.add_run(); run.text=text; run.font.size=Pt(sz); run.font.color.rgb=c; run.font.bold=bold

def vline(sl,x,y,h): R(sl,x,y,0.12,h,GOLD)
def hline(sl,x,y,w): R(sl,x,y,w,0.07,GOLD)

def page_header(sl,num,title):
    bg(sl); R(sl,0,0,0.5,14.3,DGREEN); vline(sl,0.7,1.0,12.0)
    T(sl,num,1.1,0.5,2,1.0,sz=36,c=DGREEN,bold=True)
    T(sl,title,3.0,0.5,20,1.0,sz=22,c=GOLD,bold=True)
    hline(sl,1.1,1.7,23.5)

def stat_card(sl,x,y,w,h,num,label,nc=GOLD):
    R(sl,x,y,w,h,CARD)
    T(sl,num,x,y+0.25,w,1.3,sz=28,c=nc,bold=True,align=PP_ALIGN.CENTER)
    T(sl,label,x,y+1.55,w,0.8,sz=10,c=LGRAY,align=PP_ALIGN.CENTER)

def card(sl,x,y,w,h,title,lines,tc=GOLD,lc=LGRAY,bgc=CARD):
    R(sl,x,y,w,h,bgc)
    T(sl,title,x+0.25,y+0.2,w-0.5,0.8,sz=13,c=tc,bold=True)
    tb=sl.shapes.add_textbox(Cm(x+0.25),Cm(y+1.05),Cm(w-0.5),Cm(h-1.2))
    tf=tb.text_frame; tf.word_wrap=True
    for i,line in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        run=p.add_run(); run.text=line; run.font.size=Pt(11); run.font.color.rgb=lc

# ── SLIDE 1: Cover ──────────────────────────────────────
sl=S(); bg(sl)
R(sl,0,0,0.5,14.3,DGREEN); R(sl,0,11.2,25.4,0.08,GOLD)
T(sl,'信任的杠杆',2,2.2,22,3.0,sz=54,c=GOLD,bold=True,align=PP_ALIGN.CENTER)
T(sl,'中医AI诊疗电商服务平台',2,5.5,22,1.4,sz=24,c=WHITE,align=PP_ALIGN.CENTER)
T(sl,'以AI重构中医信任链，让大师智慧触达千万家庭',2,7.0,22,1.0,sz=15,c=LGRAY,align=PP_ALIGN.CENTER,italic=True)
hline(sl,2,8.2,21.4)
for i,(n,l,nc) in enumerate([('3200亿','食药同源市场',GOLD),('24.5%','年复合增长率',GREEN),('10个月','预计回本周期',GOLD),('Pre-A轮','融资需求·500万',ORANGE)]):
    stat_card(sl,1.0+i*6.1,8.6,5.6,2.5,n,l,nc)
T(sl,'B2B2C模式重构大健康产业  ·  商业计划书  ·  2026年1月',2,12.5,22,0.7,sz=11,c=LGRAY,align=PP_ALIGN.CENTER)

# ── SLIDE 2: TOC ────────────────────────────────────────
sl=S(); bg(sl); R(sl,0,0,0.5,14.3,DGREEN)
T(sl,'目  录',1.2,0.4,22,1.0,sz=28,c=GOLD,bold=True); hline(sl,1.2,1.6,23.5)
left=[('01','执行摘要','核心价值与关键指标'),('02','行业痛点','B2C模式三大困境'),('03','解决方案','中医大师的Shopify'),('04','三步走战略','效率→依赖→变现'),('05','核心优势','零CAC + 高LTV'),('06','市场机遇','3200亿黄金赛道')]
right=[('07','竞品分析','差异化竞争定位'),('08','商业飞轮','信任驱动LTV闭环'),('09','财务模型','可预测单元经济'),('10','护城河','三重壁垒构建'),('11','风险评估','六维风险管控'),('12','执行计划','12个月里程碑')]
for i,(num,title,desc) in enumerate(left):
    y=2.1+i*1.95; R(sl,1.2,y,0.75,0.75,GOLD)
    T(sl,num,1.2,y+0.05,0.75,0.65,sz=13,c=BG,bold=True,align=PP_ALIGN.CENTER)
    T(sl,title,2.2,y,5,0.55,sz=14,c=WHITE,bold=True); T(sl,desc,2.2,y+0.6,5,0.5,sz=11,c=LGRAY)
for i,(num,title,desc) in enumerate(right):
    y=2.1+i*1.95; R(sl,13.2,y,0.75,0.75,GOLD)
    T(sl,num,13.2,y+0.05,0.75,0.65,sz=13,c=BG,bold=True,align=PP_ALIGN.CENTER)
    T(sl,title,14.2,y,5,0.55,sz=14,c=WHITE,bold=True); T(sl,desc,14.2,y+0.6,5,0.5,sz=11,c=LGRAY)

# ── SLIDE 3: Executive Summary ──────────────────────────
sl=S(); page_header(sl,'01','执行摘要')
T(sl,'我们是中医大师的"Shopify"平台，通过B2B2C模式，将医生的"信任资产"转化为可规模化的商业价值',1.1,2.0,23,0.9,sz=13,c=LGRAY,italic=True)
T(sl,'「我们不卖流量，我们卖信任」',1.1,2.9,23,0.7,sz=14,c=GOLD,bold=True)
for i,(n,l,nc) in enumerate([('3200亿','食药同源市场规模',GOLD),('24.5%','中草药年复合增长率',GREEN),('82%','愿为信任付溢价',GOLD),('10个月','预计投资回本周期',GREEN)]):
    stat_card(sl,1.1+i*6.0,3.8,5.6,2.6,n,l,nc)
hline(sl,1.1,6.7,23.5)
T(sl,'商业模式亮点',1.1,7.0,10,0.7,sz=14,c=GOLD,bold=True)
T(sl,'目标规模',13,7.0,10,0.7,sz=14,c=GOLD,bold=True)
for i,pt in enumerate(['✦ B2B2C信任赋能模式，零边际获客成本','✦ AI诊室+私域CRM+线上商城三位一体','✦ 精标药房供应链，出口级品质保障','✦ 医嘱驱动高转化，预估转化率65%','✦ 双重护城河：资源壁垒+技术壁垒']):
    T(sl,pt,1.1,7.8+i*1.1,11,0.8,sz=12,c=LGRAY)
for i,pt in enumerate(['✦ 首年：签约20名顶尖中医师','✦ 三年：覆盖200名医生，20,000名患者','✦ 盈亏平衡：50名医生时达成','✦ 三年营收目标：1800万元/年','✦ 启动A轮融资，估值目标5000万']):
    T(sl,pt,13,7.8+i*1.1,11,0.8,sz=12,c=LGRAY)

# ── SLIDE 4: Pain Points ────────────────────────────────
sl=S(); page_header(sl,'02','行业痛点分析：B2C模式的三大困境')
T(sl,'传统中医健康平台面临结构性困境，市场存在巨大的"信任赤字"亟待填补',1.1,2.0,23,0.7,sz=13,c=LGRAY,italic=True)
card(sl,1.1,3.0,7.3,9.5,'🔴  流量成本黑洞',
    ['传统B2C平台获客成本高达营收7.4%','流量属于平台，医生无法沉淀用户资产','复购率不足30%，用户忠诚度极低','每获取一名新患者成本：¥200-500','平台抽佣+广告费吞噬80%利润空间'],tc=RED)
card(sl,9.05,3.0,7.3,9.5,'🟡  信任赤字危机',
    ['市场存在数十亿级"信任赤字"缺口','82%消费者愿为信任付20-30%溢价','但现有平台只会打价格战，无法承接','假冒伪劣药材事件频发，信任崩塌','消费者迫切需要"可信赖的医生背书"'],tc=ORANGE)
card(sl,17.0,3.0,7.3,9.5,'🟢  医生IP被稀释',
    ['顶尖中医日均接诊上限仅20人','90%潜在价值被低效线下模式禁锢','面临"三流药材"绑架的风险','医生无法建立个人品牌和数字资产','核心资产"信任"无法规模化变现'],tc=GREEN)

# ── SLIDE 5: Solution ───────────────────────────────────
sl=S(); page_header(sl,'03','解决方案：中医大师的数字化操作系统')
T(sl,'我们不是又一个医疗电商平台——我们是赋能中医大师的"Shopify"，让信任成为可交易的商业资产',1.1,2.0,23,0.8,sz=13,c=LGRAY,italic=True)
for i,(icon,title,pts,c) in enumerate([
    ('🏥','AI 诊  室',['语音转文字，自动生成结构化病历','智能辨证分析，处方推荐参考','节省50%问诊时间，相当于多接诊10人/天','医案自动归档，持续训练专属AI模型'],GOLD),
    ('📱','私域 CRM',['患者健康档案数字化管理','智能随访提醒，提升复诊率','健康数据追踪，个性化调理方案','私域流量沉淀，患者资产归医生所有'],GREEN),
    ('🛒','信任商城',['医嘱一键直连购买，转化率预估65%','精标药房供应链，出口级质检标准','全程溯源码体系，药材来源透明','医生品牌背书，溢价空间20-30%'],ORANGE)]):
    x=1.1+i*8.1
    R(sl,x,3.0,7.6,8.5,CARD)
    T(sl,icon+' '+title,x+0.3,3.2,7.0,0.9,sz=16,c=c,bold=True)
    hline(sl,x+0.3,4.2,7.0)
    for j,pt in enumerate(pts):
        T(sl,'• '+pt,x+0.3,4.5+j*1.1,7.0,0.9,sz=11,c=LGRAY)
T(sl,'竞品对比：我们 vs 平安好医生 vs 京东健康',1.1,11.8,10,0.6,sz=12,c=GOLD,bold=True)
for i,(label,us,pa,jd) in enumerate([('获客成本','零边际CAC','高（广告驱动）','高（流量购买）'),('信任背书','名医个人品牌','平台品牌','平台品牌'),('医生绑定','AI深度绑定','松散合作','松散合作'),('供应链','精标药房溯源','标准化','标准化')]):
    x=1.1+i*6.0
    T(sl,label,x,12.5,5.5,0.5,sz=10,c=GOLD,bold=True,align=PP_ALIGN.CENTER)
    T(sl,us,x,13.1,5.5,0.5,sz=10,c=GREEN,align=PP_ALIGN.CENTER)
    T(sl,pa,x,13.6,5.5,0.4,sz=9,c=LGRAY,align=PP_ALIGN.CENTER)

# ── SLIDE 6: 3-Step Strategy ────────────────────────────
sl=S(); page_header(sl,'04','三步走战略：从效率工具到信任变现')
T(sl,'通过三阶段渐进式策略，将医生从"工具用户"转化为"平台共生体"',1.1,2.0,23,0.7,sz=13,c=LGRAY,italic=True)
for i,(phase,title,pts,c,period) in enumerate([
    ('第一步','效率工具',['私域CRM：纸质病历数字化','AI诊室：节省50%问诊时间','医生平均节省2小时/天','相当于多接诊6名患者/天','快速建立使用习惯与依赖'],GOLD,'M1-M3'),
    ('第二步','建立依赖',['AI学习医生医案成为数字副手','平台成为医生的数据银行','首年积累5000+专属医案','AI模型越用越懂医生','数据资产无法迁移转移'],GREEN,'M4-M6'),
    ('第三步','信任变现',['个人线上商城：医生品牌背书','精标药房：出口级可溯源药材','商城GMV目标：单医生年均36万','医嘱转化率预估65%','平台年收入：单医生3.6万'],ORANGE,'M7-M12')]):
    x=1.1+i*8.1
    R(sl,x,2.9,7.6,9.8,CARD)
    R(sl,x,2.9,7.6,0.6,c)
    T(sl,period,x+0.2,2.95,3,0.5,sz=11,c=BG,bold=True)
    T(sl,phase,x+0.2,3.7,7.2,0.7,sz=13,c=c,bold=True)
    T(sl,title,x+0.2,4.4,7.2,0.8,sz=18,c=WHITE,bold=True)
    hline(sl,x+0.2,5.3,7.0)
    for j,pt in enumerate(pts): T(sl,'▸ '+pt,x+0.2,5.6+j*1.1,7.2,0.9,sz=11,c=LGRAY)
# Arrow connectors
for i in range(2): T(sl,'→',9.3+i*8.1,7.0,1.0,1.0,sz=28,c=GOLD,bold=True,align=PP_ALIGN.CENTER)

# ── SLIDE 7: Core Advantages ────────────────────────────
sl=S(); page_header(sl,'05','核心优势：重构"流量"为"信任"')
T(sl,'我们的竞争优势不是技术领先，而是商业模式的结构性优势——信任无法被复制',1.1,2.0,23,0.7,sz=13,c=LGRAY,italic=True)
for i,(title,num,numlabel,pts,c) in enumerate([
    ('零边际获客成本','0%','vs 竞品7.4%营收',['传统B2C：CAC高达营收7.4%','我们：跳过购买流量和培养信任','本质：迁移已存在的高价值私域','竞争对手需18-24个月建立信任','我们Day1即拥有医生的信任背书'],GOLD),
    ('高转化·高LTV','65%','医嘱转化率预估',['B2C路径：曝光→比价→购买（3-5%）','我们路径：信任→医嘱→购买（65%）','核心：医嘱是终极购买理由','年人均消费：12,000元','LTV是行业平均的8-10倍'],GREEN),
    ('供应链护城河','100%','精标药房溯源',['出口级质检标准，超越国标','全程溯源码体系，来源透明','与顶级药材产地直采合作','药材品质强化医生信任背书','高品质=高溢价=高毛利（65%+）'],ORANGE)]):
    x=1.1+i*8.1
    R(sl,x,2.9,7.6,9.8,CARD)
    T(sl,num,x,3.2,7.6,2.0,sz=42,c=c,bold=True,align=PP_ALIGN.CENTER)
    T(sl,numlabel,x,5.1,7.6,0.6,sz=11,c=LGRAY,align=PP_ALIGN.CENTER)
    hline(sl,x+0.3,5.8,7.0)
    T(sl,title,x+0.3,6.1,7.0,0.8,sz=14,c=c,bold=True)
    for j,pt in enumerate(pts): T(sl,'• '+pt,x+0.3,7.0+j*1.0,7.0,0.8,sz=11,c=LGRAY)

# ── SLIDE 8: Market Opportunity ─────────────────────────
sl=S(); page_header(sl,'06','市场机遇：精准切入黄金赛道')
T(sl,'大健康产业进入爆发期，食药同源赛道是增速最快、信任溢价最高的细分市场',1.1,2.0,23,0.7,sz=13,c=LGRAY,italic=True)
for i,(n,l,desc,nc) in enumerate([('3200亿','食药同源市场(2025)','消费者从被动治疗转向主动预防，市场快速扩张',GOLD),('24.5%','中草药CAGR','市场每3年翻一番，增量红利持续释放',GREEN),('1000亿','养生茶市场(2028)','食药同源产品最佳载体，复购率高',ORANGE),('50亿','中药配方颗粒','便捷化是不可逆趋势，高毛利品类',GOLD)]):
    x=1.1+i*6.0
    R(sl,x,2.9,5.6,4.5,CARD)
    T(sl,n,x,3.1,5.6,1.8,sz=30,c=nc,bold=True,align=PP_ALIGN.CENTER)
    T(sl,l,x,4.8,5.6,0.6,sz=11,c=WHITE,bold=True,align=PP_ALIGN.CENTER)
    T(sl,desc,x+0.2,5.5,5.2,1.2,sz=10,c=LGRAY)
hline(sl,1.1,7.8,23.5)
T(sl,'目标用户画像',1.1,8.2,7,0.7,sz=14,c=GOLD,bold=True)
T(sl,'政策利好',9,8.2,7,0.7,sz=14,c=GOLD,bold=True)
T(sl,'竞争格局',17,8.2,7,0.7,sz=14,c=GOLD,bold=True)
for pt in ['年龄：35-55岁，注重健康管理','城市：一二线城市为主','收入：年收入15万元以上','特征：信任专业医生建议','痛点：缺乏可信赖的健康顾问']:
    T(sl,'• '+pt,1.1,8.9+['年龄','城市','收入','特征','痛点'].index(pt[:2])*0.9,7.5,0.8,sz=11,c=LGRAY)
for i,pt in enumerate(['国家中医药振兴计划2030','互联网医疗政策持续放开','医保数字化改革加速','中医药出海战略支持']):
    T(sl,'✓ '+pt,9,8.9+i*0.9,7.5,0.8,sz=11,c=LGRAY)
for i,pt in enumerate(['市场高度分散，CR5<15%','无主导性信任平台','存在巨大整合机会','先发优势窗口期：12-18个月']):
    T(sl,'★ '+pt,17,8.9+i*0.9,7.5,0.8,sz=11,c=LGRAY)

# ── SLIDE 9: Competitor Analysis ────────────────────────
sl=S(); page_header(sl,'07','竞品分析：差异化定位')
T(sl,'我们不与现有平台正面竞争——我们开辟了一个全新的"信任赋能"赛道',1.1,2.0,23,0.7,sz=13,c=LGRAY,italic=True)
headers=['维度','我们（信任的杠杆）','平安好医生','京东健康','好大夫在线']
widths=[3.5,5.5,4.0,4.0,4.0]
xs=[1.1,4.7,10.3,14.4,18.5]
for j,(h,w,x) in enumerate(zip(headers,widths,xs)):
    R(sl,x,2.9,w,0.8,DGREEN if j==0 else (GOLD if j==1 else MCARD))
    T(sl,h,x+0.1,2.95,w-0.2,0.7,sz=12,c=BG if j==1 else WHITE,bold=True,align=PP_ALIGN.CENTER)
rows=[
    ('商业模式','B2B2C信任赋能','B2C流量平台','B2C电商平台','B2C问诊平台'),
    ('获客成本','零边际CAC ✓','高（广告）✗','高（流量）✗','高（广告）✗'),
    ('信任来源','名医个人品牌 ✓','平台品牌 △','平台品牌 △','医生评分 △'),
    ('医生绑定','AI深度绑定 ✓','松散合作 ✗','松散合作 ✗','松散合作 ✗'),
    ('供应链','精标溯源 ✓','标准化 △','标准化 △','无 ✗'),
    ('毛利率','65%+ ✓','35-45% △','25-35% ✗','40-50% △'),
    ('核心壁垒','信任+数据双壁垒','品牌+资金','品牌+资金','品牌+流量'),
]
for i,row in enumerate(rows):
    bg_c=CARD if i%2==0 else MCARD
    for j,(val,w,x) in enumerate(zip(row,widths,xs)):
        R(sl,x,3.8+i*1.3,w,1.2,bg_c)
        c=GREEN if '✓' in val else (RED if '✗' in val else LGRAY)
        if j==1: c=GREEN if '✓' in val else (RED if '✗' in val else GOLD)
        T(sl,val,x+0.1,3.9+i*1.3,w-0.2,1.0,sz=11,c=c,align=PP_ALIGN.CENTER)
T(sl,'结论：我们是唯一以"医生信任资产"为核心的B2B2C平台，竞争对手无法在短期内复制此模式',1.1,13.0,23,0.8,sz=12,c=GOLD,bold=True)

# ── SLIDE 10: Business Flywheel ─────────────────────────
sl=S(); page_header(sl,'08','商业飞轮：信任驱动的LTV闭环')
T(sl,'每一个环节都在强化信任，信任越强，LTV越高，飞轮越快',1.1,2.0,23,0.7,sz=13,c=LGRAY,italic=True)
# Flywheel center
R(sl,9.7,5.5,6.0,3.2,GOLD)
T(sl,'信任飞轮',9.7,5.8,6.0,1.0,sz=18,c=BG,bold=True,align=PP_ALIGN.CENTER)
T(sl,'越转越快',9.7,6.8,6.0,0.8,sz=13,c=BG,align=PP_ALIGN.CENTER)
T(sl,'LTV持续增长',9.7,7.5,6.0,0.7,sz=11,c=BG,align=PP_ALIGN.CENTER)
# Flywheel nodes
nodes=[
    (1.1,2.9,'① 医生加入平台','获得AI诊室+CRM工具\n效率提升50%',GOLD),
    (17.5,2.9,'② AI深度学习','积累专属医案数据\n形成数字副手',GREEN),
    (17.5,8.5,'③ 信任商城上线','医嘱直连购买\n转化率65%',ORANGE),
    (1.1,8.5,'④ 患者复购增长','LTV年均12,000元\n口碑裂变新患者',GOLD),
]
for x,y,title,desc,c in nodes:
    R(sl,x,y,6.2,3.8,CARD)
    T(sl,title,x+0.2,y+0.3,5.8,0.8,sz=13,c=c,bold=True)
    T(sl,desc,x+0.2,y+1.2,5.8,1.8,sz=11,c=LGRAY)
# Arrows
for x,y,txt in [(7.5,4.5,'→'),(15.5,4.5,'↘'),(15.5,9.5,'←'),(7.5,9.5,'↙')]:
    T(sl,txt,x,y,1.5,1.0,sz=24,c=GOLD,bold=True,align=PP_ALIGN.CENTER)
# KPIs
hline(sl,1.1,12.8,23.5)
for i,(n,l) in enumerate([('65%','医嘱转化率'),('12,000元','年人均LTV'),('8-10x','vs行业平均LTV'),('36万','单医生年GMV')]):
    T(sl,n,1.1+i*6.0,13.1,5.6,0.8,sz=22,c=GOLD,bold=True,align=PP_ALIGN.CENTER)
    T(sl,l,1.1+i*6.0,13.9,5.6,0.5,sz=10,c=LGRAY,align=PP_ALIGN.CENTER)

# ── SLIDE 11: Financial Model ────────────────────────────
sl=S(); page_header(sl,'09','财务模型：可预测的单元经济')
T(sl,'基于保守假设的单元经济模型，展示清晰的盈利路径',1.1,2.0,23,0.7,sz=13,c=LGRAY,italic=True)
# Unit economics
for i,(title,rows2,c) in enumerate([
    ('单医生单元经济（年）',[('患者数量','200名活跃患者'),('年人均消费','¥12,000'),('商城GMV','¥240万'),('平台抽佣(15%)','¥36万'),('毛利率','65%+')],GOLD),
    ('平台规模经济（3年）',[('医生数量','200名签约医生'),('总GMV','¥4.8亿'),('平台年收入','¥7,200万'),('净利润率','35%'),('净利润','¥2,520万')],GREEN),
    ('融资与回报',[('本轮融资','Pre-A 500万'),('估值','2,500万（5x）'),('回本周期','10个月'),('A轮目标','5,000万估值'),('IRR预期','300%+')],ORANGE)]):
    x=1.1+i*8.1
    R(sl,x,2.9,7.6,9.8,CARD)
    T(sl,title,x+0.2,3.1,7.2,0.8,sz=13,c=c,bold=True)
    hline(sl,x+0.2,4.0,7.0)
    for j,(k,v) in enumerate(rows2):
        R(sl,x+0.2,4.3+j*1.6,7.0,1.4,MCARD if j%2==0 else CARD)
        T(sl,k,x+0.4,4.4+j*1.6,3.5,0.6,sz=11,c=LGRAY)
        T(sl,v,x+0.4,5.0+j*1.6,6.6,0.6,sz=12,c=c,bold=True)
hline(sl,1.1,13.0,23.5)
T(sl,'盈亏平衡点：签约50名医生时达成  |  首年目标：20名医生  |  三年目标：200名医生',1.1,13.2,23,0.7,sz=12,c=GOLD,bold=True,align=PP_ALIGN.CENTER)

# ── SLIDE 12: Moat ───────────────────────────────────────
sl=S(); page_header(sl,'10','护城河：三重壁垒构建')
T(sl,'我们的竞争壁垒不依赖资金，而是依赖时间积累——越早建立，越难被复制',1.1,2.0,23,0.7,sz=13,c=LGRAY,italic=True)
for i,(num,title,pts,c) in enumerate([
    ('壁垒一','资源壁垒（立即生效）',['顶尖中医师资源极度稀缺','全国执业中医师约80万人','三甲医院主任医师仅约2万人','我们目标签约其中TOP 200','一旦签约，竞争对手无法挖角'],GOLD),
    ('壁垒二','数据壁垒（6个月建立）',['AI持续学习每位医生的医案','专属模型越用越懂医生习惯','数据资产归医生所有但存于平台','迁移成本极高（数年积累）','形成"数字副手"强依赖关系'],GREEN),
    ('壁垒三','信任壁垒（12个月固化）',['患者信任绑定在医生个人品牌','医生品牌与平台深度融合','患者不会因平台切换而流失','口碑裂变形成自增长飞轮','信任网络一旦形成无法复制'],ORANGE)]):
    x=1.1+i*8.1
    R(sl,x,2.9,7.6,9.8,CARD)
    R(sl,x,2.9,7.6,0.5,c)
    T(sl,num,x+0.2,3.0,7.2,0.4,sz=10,c=BG,bold=True)
    T(sl,title,x+0.2,3.6,7.2,0.8,sz=13,c=c,bold=True)
    hline(sl,x+0.2,4.5,7.0)
    for j,pt in enumerate(pts): T(sl,'▸ '+pt,x+0.2,4.8+j*1.1,7.2,0.9,sz=11,c=LGRAY)
T(sl,'综合评估：三重壁垒叠加，形成"时间+资源+信任"的复合护城河，竞争对手需3-5年才能追赶',1.1,13.0,23,0.8,sz=12,c=GOLD,bold=True)

# ── SLIDE 13: Execution Plan ─────────────────────────────
sl=S(); page_header(sl,'11','执行计划：12个月里程碑')
T(sl,'清晰的执行路径，每个阶段都有可量化的成功指标',1.1,2.0,23,0.7,sz=13,c=LGRAY,italic=True)
milestones=[
    ('M1-M2','产品打磨',GOLD,['完成AI诊室MVP','招募首批5名种子医生','完成产品迭代反馈循环','建立基础CRM功能']),
    ('M3-M4','验证模型',GREEN,['签约10名医生','首批商城上线','验证65%医嘱转化率','收集用户反馈优化']),
    ('M5-M6','规模复制',ORANGE,['签约20名医生','月GMV突破100万','精标药房供应链就绪','完成Pre-A融资']),
    ('M7-M9','加速扩张',GOLD,['签约50名医生','月GMV突破500万','达到盈亏平衡点','启动A轮融资准备']),
    ('M10-M12','生态建设',GREEN,['签约100名医生','年GMV突破1亿','建立医生社区生态','完成A轮融资']),
]
for i,(period,phase,c,pts) in enumerate(milestones):
    x=1.1+i*4.7
    R(sl,x,2.9,4.3,10.0,CARD)
    R(sl,x,2.9,4.3,0.5,c)
    T(sl,period,x+0.1,3.0,4.1,0.4,sz=10,c=BG,bold=True,align=PP_ALIGN.CENTER)
    T(sl,phase,x+0.1,3.6,4.1,0.7,sz=13,c=c,bold=True,align=PP_ALIGN.CENTER)
    hline(sl,x+0.2,4.4,3.8)
    for j,pt in enumerate(pts): T(sl,'• '+pt,x+0.2,4.7+j*1.5,3.9,1.2,sz=10,c=LGRAY)
# Timeline bar
R(sl,1.1,13.1,23.5,0.5,MCARD)
for i,c in enumerate([GOLD,GREEN,ORANGE,GOLD,GREEN]):
    R(sl,1.1+i*4.7,13.1,4.3,0.5,c)
    T(sl,milestones[i][0],1.2+i*4.7,13.15,4.1,0.4,sz=10,c=BG,bold=True,align=PP_ALIGN.CENTER)

# ── SLIDE 14: Closing / CTA ──────────────────────────────
sl=S()
R(sl,0,0,25.4,19.05,BG)
# Gradient accent
R(sl,0,0,25.4,0.4,GOLD)
R(sl,0,18.65,25.4,0.4,GOLD)
T(sl,'现在，是加入的最佳时机',1.1,1.5,23,1.5,sz=36,c=WHITE,bold=True,align=PP_ALIGN.CENTER)
T(sl,'先发优势窗口期只有12-18个月',1.1,3.2,23,0.8,sz=18,c=GOLD,align=PP_ALIGN.CENTER)
hline(sl,4.0,4.3,17.5)
T(sl,'我们正在寻找：',1.1,5.0,23,0.7,sz=14,c=LGRAY,align=PP_ALIGN.CENTER)
for i,(icon,txt) in enumerate([('💰','Pre-A轮战略投资人（500万）'),('🏥','顶尖中医师合作伙伴（首批20名）'),('🌿','出口级药材供应商（精标认证）')]):
    R(sl,4.5+i*5.5,5.9,5.0,4.5,CARD)
    T(sl,icon,4.5+i*5.5,6.2,5.0,1.2,sz=28,c=WHITE,align=PP_ALIGN.CENTER)
    T(sl,txt,4.5+i*5.5,7.5,5.0,1.5,sz=12,c=WHITE,bold=True,align=PP_ALIGN.CENTER)
hline(sl,4.0,11.0,17.5)
T(sl,'信任的杠杆  ·  中医AI平台',1.1,11.5,23,0.8,sz=20,c=GOLD,bold=True,align=PP_ALIGN.CENTER)
T(sl,'让每一位名医的信任，成为可持续增长的商业资产',1.1,12.5,23,0.7,sz=14,c=LGRAY,align=PP_ALIGN.CENTER)
for i,(k,v) in enumerate([('联系方式','contact@tcm-ai.com'),('官网','www.tcm-ai.com'),('微信','TCM_AI_Platform')]):
    T(sl,f'{k}：{v}',1.1+i*8.0,14.5,7.5,0.6,sz=12,c=LGRAY,align=PP_ALIGN.CENTER)

# ── SAVE ─────────────────────────────────────────────────
out='E:/yw/agiatme/goose/output/信任的杠杆_中医AI平台_增强版.pptx'
prs.save(out)
print(f'Saved: {out}  ({len(prs.slides)} slides)')
