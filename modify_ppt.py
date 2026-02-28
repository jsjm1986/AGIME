# -*- coding: utf-8 -*-
from pptx import Presentation

SRC = './data/workspaces/698616a1980c003c66f6421e/missions/40204152-de4f-41ac-bca3-17265a7f1112/documents/doc_b75f92060612ba80.pptx'
DST = './data/workspaces/698616a1980c003c66f6421e/missions/40204152-de4f-41ac-bca3-17265a7f1112/documents/doc_b75f92060612ba80.pptx'

prs = Presentation(SRC)

def set_text(slide, shape_idx, para_idx, new_text):
    """Replace text in a specific paragraph while preserving formatting."""
    para = slide.shapes[shape_idx].text_frame.paragraphs[para_idx]
    if para.runs:
        para.runs[0].text = new_text
        for run in para.runs[1:]:
            run.text = ''
    else:
        para.text = new_text

# ============ Slide 1: Cover - Add AI keywords ============
s1 = prs.slides[0]
set_text(s1, 4, 0, '中医AI智能诊疗电商服务平台')
set_text(s1, 5, 0, 'AI驱动的B2B2C模式 · 重构大健康产业信任链')

# ============ Slide 3: Executive Summary - Strengthen AI narrative ============
s3 = prs.slides[2]
set_text(s3, 6, 0, '我们是中医大师的"AI赋能Shopify"平台，以AI大模型+知识图谱为技术引擎，通过B2B2C模式，将医生的"信任资产"转化为可规模化的商业价值。')
set_text(s3, 25, 0, '• B2B2C信任赋能 + AI智能中枢')
set_text(s3, 26, 0, '• AI大模型诊室+私域CRM+线上商城')
set_text(s3, 29, 0, '• AI数据飞轮 + 资源双护城河')

# ============ Slide 6: Solution 三步走 - Strengthen AI in each step ============
s6 = prs.slides[5]
set_text(s6, 8, 0, '• AI智能CRM：病历数字化+智能随访')
set_text(s6, 9, 0, '• AI大模型诊室：节约50%问诊时间')
set_text(s6, 14, 0, '• AI深度学习医案，成为大师数字分身')
set_text(s6, 15, 0, '• 知识图谱沉淀，平台成为AI数据银行')

# ============ Slide 7: Core Advantages - Add AI tech advantage ============
s7 = prs.slides[6]
set_text(s7, 3, 0, '核心优势：AI重构"流量"与"信任"')
set_text(s7, 8, 0, '• 我们: AI精准匹配，跳过流量和信任培养')
set_text(s7, 14, 0, '• 我们: 信任→AI医嘱→购买(智能强转化)')
set_text(s7, 15, 0, '• 核心: AI辅助医嘱是终极购买理由')

# ============ Slide 8: Market Opportunity - Add AI healthcare data ============
s8 = prs.slides[7]
set_text(s8, 3, 0, '市场机遇：AI+大健康黄金赛道')
set_text(s8, 8, 0, '消费者从被动治疗转向AI主动健康管理')
set_text(s8, 13, 0, 'AI+中医赛道，市场每3年翻一番')
set_text(s8, 18, 0, 'AI个性化推荐驱动的最佳消费载体')
set_text(s8, 23, 0, 'AI智能配方+便捷化是不可逆趋势')

# ============ Slide 11: Moat Strategy - Strengthen AI flywheel ============
s11 = prs.slides[10]
set_text(s11, 3, 0, '护城河策略：资源+AI数据飞轮双壁垒')
set_text(s11, 14, 0, '• B2C巨头只有通用AI，无名医数据')
set_text(s11, 15, 0, '• 我们：每位大师专属AI模型(数字分身)')
set_text(s11, 16, 0, '• 医案数据越多→AI越精准→医生越依赖')
set_text(s11, 17, 0, '• 数据飞轮效应：AI壁垒随时间指数增长')

prs.save(DST)
print('PPT modified successfully!')
