from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# ── Color Palette ──────────────────────────────────────────────
C_BG_DARK   = RGBColor(0x0A, 0x0E, 0x27)   # deep navy
C_BG_MID    = RGBColor(0x0F, 0x17, 0x3A)   # mid navy
C_ACCENT1   = RGBColor(0x00, 0xC2, 0xFF)   # electric cyan
C_ACCENT2   = RGBColor(0xF5, 0xC5, 0x18)   # gold
C_ACCENT3   = RGBColor(0x7B, 0x2F, 0xFF)   # purple
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT     = RGBColor(0xCC, 0xDD, 0xFF)   # soft blue-white
C_GRAY      = RGBColor(0x88, 0x99, 0xBB)
C_CARD      = RGBColor(0x14, 0x1E, 0x4A)   # card bg
C_GREEN     = RGBColor(0x00, 0xE5, 0x96)   # mint green

W, H = Inches(13.33), Inches(7.5)   # 16:9

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

# ── Helpers ────────────────────────────────────────────────────
def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)

def fill_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE=1
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             font_size=Pt(18), bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb

def add_multiline(slide, lines, l, t, w, h,
                  font_size=Pt(16), bold=False, color=C_WHITE,
                  align=PP_ALIGN.LEFT, line_spacing=None):
    """lines: list of (text, size, bold, color) or plain str"""
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            txt, sz, bd, cl = item, font_size, bold, color
        else:
            txt, sz, bd, cl = item
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = txt
        run.font.size = sz
        run.font.bold = bd
        run.font.color.rgb = cl
    return txb

def gradient_rect(slide, l, t, w, h, color1, color2, angle=0):
    """Approximate gradient with two overlapping rects (pptx limitation)."""
    r = add_rect(slide, l, t, w, h, fill_color=color1)
    return r

def add_circle(slide, cx, cy, r_inch, fill_color, line_color=None):
    from pptx.util import Inches
    l = cx - r_inch/2
    t = cy - r_inch/2
    shape = slide.shapes.add_shape(9, l, t, r_inch, r_inch)  # oval
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

print("Helpers loaded.")

# ══════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════
s1 = blank_slide(prs)
fill_bg(s1, C_BG_DARK)

# Left accent bar
add_rect(s1, Inches(0), Inches(0), Inches(0.08), H, fill_color=C_ACCENT1)

# Decorative large circle (background)
add_circle(s1, Inches(10.5), Inches(3.75), Inches(6.5),
           RGBColor(0x0F, 0x1A, 0x45))
add_circle(s1, Inches(10.8), Inches(3.5), Inches(4.8),
           RGBColor(0x12, 0x20, 0x55))
add_circle(s1, Inches(11.0), Inches(3.3), Inches(3.0),
           RGBColor(0x00, 0xC2, 0xFF), line_color=C_ACCENT1)

# Tag line bar
add_rect(s1, Inches(0.5), Inches(1.2), Inches(3.2), Inches(0.38),
         fill_color=C_ACCENT1)
add_text(s1, "AI-Native 出版协作平台",
         Inches(0.55), Inches(1.22), Inches(3.1), Inches(0.34),
         font_size=Pt(13), bold=True, color=C_BG_DARK, align=PP_ALIGN.LEFT)

# Main title
add_text(s1, "编得好",
         Inches(0.5), Inches(1.75), Inches(7), Inches(1.5),
         font_size=Pt(72), bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)

# Subtitle
add_text(s1, "全国出版编校智能协作网络",
         Inches(0.5), Inches(3.1), Inches(8), Inches(0.7),
         font_size=Pt(28), bold=False, color=C_ACCENT1, align=PP_ALIGN.LEFT)

# Description
add_text(s1, "重塑出版行业生产关系  · 打造亿级垂直数据资产",
         Inches(0.5), Inches(3.85), Inches(8.5), Inches(0.5),
         font_size=Pt(16), bold=False, color=C_LIGHT, align=PP_ALIGN.LEFT)

# Divider line
add_rect(s1, Inches(0.5), Inches(4.5), Inches(5), Inches(0.03),
         fill_color=C_ACCENT2)

# Bottom info
add_text(s1, "项目立项建议书  ·  投资人专版",
         Inches(0.5), Inches(4.65), Inches(6), Inches(0.4),
         font_size=Pt(13), color=C_GRAY, align=PP_ALIGN.LEFT)

# Right side — 3 key stats
stats = [
    ("10×", "交易效率提升"),
    ("50%", "出版周期缩短"),
    ("亿级", "数据资产估值"),
]
for i, (num, label) in enumerate(stats):
    bx = Inches(9.2)
    by = Inches(1.5 + i * 1.7)
    add_rect(s1, bx, by, Inches(3.5), Inches(1.3),
             fill_color=C_CARD, line_color=C_ACCENT1, line_width=Pt(1))
    add_text(s1, num, bx + Inches(0.15), by + Inches(0.05),
             Inches(3.2), Inches(0.75),
             font_size=Pt(38), bold=True, color=C_ACCENT2, align=PP_ALIGN.CENTER)
    add_text(s1, label, bx + Inches(0.15), by + Inches(0.78),
             Inches(3.2), Inches(0.4),
             font_size=Pt(13), color=C_LIGHT, align=PP_ALIGN.CENTER)

print("Slide 1 done.")

# ══════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA / TABLE OF CONTENTS
# ══════════════════════════════════════════════════════════════
s2 = blank_slide(prs)
fill_bg(s2, C_BG_DARK)
add_rect(s2, Inches(0), Inches(0), Inches(0.08), H, fill_color=C_ACCENT2)

# Title
add_text(s2, "目录", Inches(0.5), Inches(0.35), Inches(4), Inches(0.7),
         font_size=Pt(36), bold=True, color=C_WHITE)
add_rect(s2, Inches(0.5), Inches(1.05), Inches(1.2), Inches(0.04), fill_color=C_ACCENT2)

items = [
    ("01", "行业现状与痛点", "出版行业的三大核心挑战"),
    ("02", "核心解决方案", "AI 三大智能体群重构生产关系"),
    ("03", "平台功能全景", "四大 AI 引擎逐一击破瓶颈"),
    ("04", "实施路线图", "三阶段 12 个月落地计划"),
    ("05", "预期效益", "社会价值 + 经济回报"),
    ("06", "投资亮点", "为什么现在，为什么我们"),
]
for i, (num, title, sub) in enumerate(items):
    row = i % 3
    col = i // 3
    bx = Inches(0.5 + col * 6.4)
    by = Inches(1.4 + row * 1.8)
    add_rect(s2, bx, by, Inches(5.9), Inches(1.5),
             fill_color=C_CARD, line_color=C_ACCENT1, line_width=Pt(0.8))
    # Number badge
    add_rect(s2, bx, by, Inches(0.7), Inches(1.5), fill_color=C_ACCENT1)
    add_text(s2, num, bx + Inches(0.02), by + Inches(0.35),
             Inches(0.66), Inches(0.7),
             font_size=Pt(22), bold=True, color=C_BG_DARK, align=PP_ALIGN.CENTER)
    add_text(s2, title, bx + Inches(0.8), by + Inches(0.2),
             Inches(5.0), Inches(0.55),
             font_size=Pt(17), bold=True, color=C_WHITE)
    add_text(s2, sub, bx + Inches(0.8), by + Inches(0.78),
             Inches(5.0), Inches(0.55),
             font_size=Pt(12), color=C_GRAY)

print("Slide 2 done.")

# ══════════════════════════════════════════════════════════════
# SLIDE 3 — PAIN POINTS
# ══════════════════════════════════════════════════════════════
s3 = blank_slide(prs)
fill_bg(s3, C_BG_DARK)
add_rect(s3, Inches(0), Inches(0), Inches(0.08), H, fill_color=RGBColor(0xFF,0x4D,0x4D))

# Section label
add_rect(s3, Inches(0.5), Inches(0.3), Inches(0.9), Inches(0.35), fill_color=RGBColor(0xFF,0x4D,0x4D))
add_text(s3, "01", Inches(0.5), Inches(0.3), Inches(0.9), Inches(0.35),
         font_size=Pt(14), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(s3, "行业现状与痛点", Inches(1.5), Inches(0.3), Inches(6), Inches(0.45),
         font_size=Pt(26), bold=True, color=C_WHITE)
add_rect(s3, Inches(0.5), Inches(0.82), Inches(12.3), Inches(0.03), fill_color=C_CARD)

# Sub-headline
add_text(s3, "出版行业正面临三大结构性挑战，传统模式已无法应对",
         Inches(0.5), Inches(0.95), Inches(12), Inches(0.4),
         font_size=Pt(14), color=C_GRAY)

pain_points = [
    (RGBColor(0xFF,0x6B,0x6B), "⚡", "供需信息不对称",
     "编校资源分布不均", 
     "传统平台靠人工搜索，效率低下\n跨地域、跨领域精准撮合几乎不可能\n优质编辑长期处于信息孤岛"),
    (RGBColor(0xFF,0xA5,0x00), "⏱", "效率低 · 流程繁琐",
     "传统接单方式严重落后",
     "依靠QQ/微信发包，缺乏工具支撑\n合同、发票、个税手动处理耗时耗力\n沟通成本高，出版周期难以压缩"),
    (RGBColor(0xFF,0x4D,0xA6), "🔍", "质量评价缺失",
     "信用体系空白，劣币驱逐良币",
     "评价主观性强，难以服众\n无法客观量化编辑贡献与质量\n优秀编辑无法获得应有回报"),
]
for i, (accent, icon, title, sub, detail) in enumerate(pain_points):
    bx = Inches(0.4 + i * 4.3)
    by = Inches(1.5)
    bw = Inches(4.0)
    bh = Inches(5.5)
    # Card
    add_rect(s3, bx, by, bw, bh, fill_color=C_CARD, line_color=accent, line_width=Pt(1.5))
    # Top accent bar
    add_rect(s3, bx, by, bw, Inches(0.08), fill_color=accent)
    # Icon circle
    add_circle(s3, bx + Inches(0.55), by + Inches(0.65), Inches(0.7),
               RGBColor(0x1A, 0x25, 0x55), line_color=accent)
    add_text(s3, icon, bx + Inches(0.2), by + Inches(0.3),
             Inches(0.7), Inches(0.6),
             font_size=Pt(22), align=PP_ALIGN.CENTER, color=accent)
    add_text(s3, title, bx + Inches(0.15), by + Inches(1.0),
             Inches(3.7), Inches(0.55),
             font_size=Pt(18), bold=True, color=C_WHITE)
    add_text(s3, sub, bx + Inches(0.15), by + Inches(1.55),
             Inches(3.7), Inches(0.4),
             font_size=Pt(12), color=accent)
    add_rect(s3, bx + Inches(0.15), by + Inches(2.0), Inches(3.5), Inches(0.03), fill_color=accent)
    add_text(s3, detail, bx + Inches(0.15), by + Inches(2.15),
             Inches(3.7), Inches(2.8),
             font_size=Pt(13), color=C_LIGHT)

# Bottom call-out
add_rect(s3, Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.38),
         fill_color=RGBColor(0x1A, 0x10, 0x10), line_color=RGBColor(0xFF,0x4D,0x4D), line_width=Pt(0.8))
add_text(s3, "💡  这三大痛点，正是「编得好」AI 平台的核心突破口",
         Inches(0.6), Inches(7.02), Inches(12), Inches(0.35),
         font_size=Pt(13), bold=True, color=RGBColor(0xFF,0x9A,0x9A), align=PP_ALIGN.CENTER)

print("Slide 3 done.")

# ══════════════════════════════════════════════════════════════
# SLIDE 4 — CORE SOLUTION: AI三大智能体
# ══════════════════════════════════════════════════════════════
s4 = blank_slide(prs)
fill_bg(s4, C_BG_DARK)
add_rect(s4, Inches(0), Inches(0), Inches(0.08), H, fill_color=C_ACCENT1)

add_rect(s4, Inches(0.5), Inches(0.3), Inches(0.9), Inches(0.35), fill_color=C_ACCENT1)
add_text(s4, "02", Inches(0.5), Inches(0.3), Inches(0.9), Inches(0.35),
         font_size=Pt(14), bold=True, color=C_BG_DARK, align=PP_ALIGN.CENTER)
add_text(s4, "核心解决方案", Inches(1.5), Inches(0.3), Inches(8), Inches(0.45),
         font_size=Pt(26), bold=True, color=C_WHITE)
add_rect(s4, Inches(0.5), Inches(0.82), Inches(12.3), Inches(0.03), fill_color=C_CARD)
add_text(s4, "「1+3+N」智能体系 — AI 经纪人 + AI 助理 + AI 审计员，重构出版生产关系",
         Inches(0.5), Inches(0.95), Inches(12.5), Inches(0.4),
         font_size=Pt(14), color=C_GRAY)

# Center hub
add_circle(s4, Inches(6.67), Inches(4.2), Inches(1.8),
           RGBColor(0x0A, 0x14, 0x35), line_color=C_ACCENT1)
add_circle(s4, Inches(6.67), Inches(4.2), Inches(1.4),
           RGBColor(0x00, 0xC2, 0xFF))
add_text(s4, "多模型\n调度网关", Inches(5.97), Inches(3.65), Inches(1.4), Inches(1.1),
         font_size=Pt(12), bold=True, color=C_BG_DARK, align=PP_ALIGN.CENTER)

agents = [
    (C_ACCENT2, Inches(1.0), Inches(1.5), "发布侧 Agent",
     "出版社的\n「项目指挥官」",
     ["零表单意图识别发包", "需求拆解 & 预算风控", "智能档期调度", "语义共鸣撮合引擎"]),
    (C_GREEN, Inches(1.0), Inches(5.5), "生产侧 Agent",
     "编校者的\n「超级副驾驶」",
     ["AI Copilot 全程伴随", "智能审稿与逻辑体检", "极速阅读与总结", "稿件透视问答 ChatDOC"]),
    (C_ACCENT3, Inches(9.5), Inches(3.5), "管理侧 Agent",
     "平台的\n「公正审计员」",
     ["全样本数据审计", "价值导向动态结算", "数字信用档案", "可信职业画像系统"]),
]
for accent, bx, by, title, sub, features in agents:
    bw, bh = Inches(3.5), Inches(2.8)
    add_rect(s4, bx, by, bw, bh, fill_color=C_CARD, line_color=accent, line_width=Pt(1.5))
    add_rect(s4, bx, by, bw, Inches(0.07), fill_color=accent)
    add_text(s4, title, bx + Inches(0.15), by + Inches(0.12),
             Inches(3.2), Inches(0.45),
             font_size=Pt(17), bold=True, color=accent)
    add_text(s4, sub, bx + Inches(0.15), by + Inches(0.55),
             Inches(3.2), Inches(0.55),
             font_size=Pt(12), color=C_LIGHT)
    add_rect(s4, bx + Inches(0.15), by + Inches(1.1), Inches(3.0), Inches(0.03), fill_color=accent)
    feat_text = "\n".join(f"▸  {f}" for f in features)
    add_text(s4, feat_text, bx + Inches(0.15), by + Inches(1.2),
             Inches(3.2), Inches(1.5),
             font_size=Pt(11), color=C_LIGHT)

# Connecting lines (visual arrows as thin rects)
add_rect(s4, Inches(4.5), Inches(2.85), Inches(1.8), Inches(0.04), fill_color=C_ACCENT2)
add_rect(s4, Inches(4.5), Inches(6.0), Inches(1.8), Inches(0.04), fill_color=C_GREEN)
add_rect(s4, Inches(9.0), Inches(4.2), Inches(0.5), Inches(0.04), fill_color=C_ACCENT3)

# Bottom bar
add_rect(s4, Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.38),
         fill_color=RGBColor(0x0A, 0x18, 0x35), line_color=C_ACCENT1, line_width=Pt(0.8))
add_text(s4, "🔒  混合部署策略：公有云 API 降本 + 本地私有化部署，敏感数据绝不出域",
         Inches(0.6), Inches(7.02), Inches(12), Inches(0.35),
         font_size=Pt(12), bold=True, color=C_ACCENT1, align=PP_ALIGN.CENTER)

print("Slide 4 done.")

# ══════════════════════════════════════════════════════════════
# SLIDE 5 — PLATFORM FEATURES (四大AI引擎)
# ══════════════════════════════════════════════════════════════
s5 = blank_slide(prs)
fill_bg(s5, C_BG_DARK)
add_rect(s5, Inches(0), Inches(0), Inches(0.08), H, fill_color=C_GREEN)

add_rect(s5, Inches(0.5), Inches(0.3), Inches(0.9), Inches(0.35), fill_color=C_GREEN)
add_text(s5, "03", Inches(0.5), Inches(0.3), Inches(0.9), Inches(0.35),
         font_size=Pt(14), bold=True, color=C_BG_DARK, align=PP_ALIGN.CENTER)
add_text(s5, "平台功能全景", Inches(1.5), Inches(0.3), Inches(8), Inches(0.45),
         font_size=Pt(26), bold=True, color=C_WHITE)
add_rect(s5, Inches(0.5), Inches(0.82), Inches(12.3), Inches(0.03), fill_color=C_CARD)
add_text(s5, "四大 AI 引擎，逐一击破传统出版瓶颈",
         Inches(0.5), Inches(0.95), Inches(12), Inches(0.4),
         font_size=Pt(14), color=C_GRAY)

features = [
    (C_ACCENT1, "语义共鸣撮合引擎",
     "从「搜人」到「算人」",
     ["深度语义分析稿件学科谱系", "亿级人才库特征向量匹配", "跨域冷启动推荐稀缺专家", "动态档期实时调度"]),
    (C_ACCENT2, "云端 AI 协同工作台",
     "Zero-Install 沉浸式编辑",
     ["完美解析 Word/PDF/LaTeX", "AI Copilot 逻辑体检 & 纠错", "百万字 1 分钟生成大纲", "ChatDOC 跨章节问答"]),
    (C_GREEN, "数字信用审计网络",
     "全样本数据审计",
     ["记录每一次修改的价值", "有效修改率客观量化", "价值导向动态结算模型", "加密不可篡改信用档案"]),
    (C_ACCENT3, "AI 赋能学院",
     "降低门槛，扩大人才供给",
     ["新手导航教学模式", "AI 导师分步提示规范", "脱敏真实稿件模拟演练", "缩短成才周期 40%"]),
]
for i, (accent, title, sub, pts) in enumerate(features):
    col = i % 2
    row = i // 2
    bx = Inches(0.4 + col * 6.4)
    by = Inches(1.4 + row * 2.9)
    bw, bh = Inches(6.0), Inches(2.6)
    add_rect(s5, bx, by, bw, bh, fill_color=C_CARD, line_color=accent, line_width=Pt(1.2))
    add_rect(s5, bx, by, Inches(0.07), bh, fill_color=accent)
    add_text(s5, title, bx + Inches(0.2), by + Inches(0.1),
             Inches(5.6), Inches(0.5),
             font_size=Pt(17), bold=True, color=accent)
    add_text(s5, sub, bx + Inches(0.2), by + Inches(0.6),
             Inches(5.6), Inches(0.35),
             font_size=Pt(12), color=C_GRAY, italic=True)
    add_rect(s5, bx + Inches(0.2), by + Inches(0.98), Inches(5.5), Inches(0.03), fill_color=accent)
    pts_text = "   ".join(f"✦ {p}" for p in pts[:2]) + "\n   ".join(["", f"✦ {pts[2]}", f"✦ {pts[3]}"])
    add_text(s5, pts_text, bx + Inches(0.2), by + Inches(1.08),
             Inches(5.6), Inches(1.4),
             font_size=Pt(11.5), color=C_LIGHT)

print("Slide 5 done.")

# ══════════════════════════════════════════════════════════════
# SLIDE 6 — IMPLEMENTATION ROADMAP
# ══════════════════════════════════════════════════════════════
s6 = blank_slide(prs)
fill_bg(s6, C_BG_DARK)
add_rect(s6, Inches(0), Inches(0), Inches(0.08), H, fill_color=C_ACCENT2)

add_rect(s6, Inches(0.5), Inches(0.3), Inches(0.9), Inches(0.35), fill_color=C_ACCENT2)
add_text(s6, "04", Inches(0.5), Inches(0.3), Inches(0.9), Inches(0.35),
         font_size=Pt(14), bold=True, color=C_BG_DARK, align=PP_ALIGN.CENTER)
add_text(s6, "实施路线图", Inches(1.5), Inches(0.3), Inches(8), Inches(0.45),
         font_size=Pt(26), bold=True, color=C_WHITE)
add_rect(s6, Inches(0.5), Inches(0.82), Inches(12.3), Inches(0.03), fill_color=C_CARD)
add_text(s6, "三阶段 · 12 个月 · 从智能基座到全国生态扩展",
         Inches(0.5), Inches(0.95), Inches(12), Inches(0.4),
         font_size=Pt(14), color=C_GRAY)

phases = [
    (C_ACCENT1, "Phase 1", "智能基座", "0 - 3 个月",
     ["完成平台基础架构搭建", "上线 AI 预审工具（敏感词/校对）", "建立编校人才画像库", "稿件语义指纹算法研发", "内部试用，沉淀数据"]),
    (C_ACCENT2, "Phase 2", "人机协同", "4 - 6 个月",
     ["推出 AI Copilot 工作台", "选取 3-5 家头部出版社试点", "跑通「发包-AI辅助-验收」闭环", "数字信用审计系统上线", "收集反馈，快速迭代"]),
    (C_GREEN, "Phase 3", "生态扩展", "7 - 12 个月",
     ["开放 API 接入第三方工具", "集成主流大模型能力", "沉淀垂直领域知识库", "全面推广至全国出版单位", "构建亿级数据资产"]),
]

# Timeline bar
add_rect(s6, Inches(0.5), Inches(3.9), Inches(12.3), Inches(0.08), fill_color=C_CARD)
for i in range(3):
    cx = Inches(0.5 + i * 4.1 + 2.05)
    add_circle(s6, cx, Inches(3.94), Inches(0.35), phases[i][0])

for i, (accent, phase, title, period, tasks) in enumerate(phases):
    bx = Inches(0.4 + i * 4.3)
    # Upper card (period + title)
    add_rect(s6, bx, Inches(1.5), Inches(4.0), Inches(2.2),
             fill_color=C_CARD, line_color=accent, line_width=Pt(1.2))
    add_rect(s6, bx, Inches(1.5), Inches(4.0), Inches(0.07), fill_color=accent)
    add_text(s6, phase, bx + Inches(0.15), Inches(1.6),
             Inches(3.7), Inches(0.4),
             font_size=Pt(12), color=accent, bold=True)
    add_text(s6, title, bx + Inches(0.15), Inches(2.0),
             Inches(3.7), Inches(0.55),
             font_size=Pt(22), bold=True, color=C_WHITE)
    add_text(s6, period, bx + Inches(0.15), Inches(2.55),
             Inches(3.7), Inches(0.4),
             font_size=Pt(13), color=C_GRAY)
    # Lower card (tasks)
    add_rect(s6, bx, Inches(4.2), Inches(4.0), Inches(2.9),
             fill_color=RGBColor(0x10, 0x18, 0x38), line_color=accent, line_width=Pt(0.8))
    task_text = "\n".join(f"▸  {t}" for t in tasks)
    add_text(s6, task_text, bx + Inches(0.15), Inches(4.3),
             Inches(3.7), Inches(2.7),
             font_size=Pt(11.5), color=C_LIGHT)

print("Slide 6 done.")

# ══════════════════════════════════════════════════════════════
# SLIDE 7 — EXPECTED BENEFITS
# ══════════════════════════════════════════════════════════════
s7 = blank_slide(prs)
fill_bg(s7, C_BG_DARK)
add_rect(s7, Inches(0), Inches(0), Inches(0.08), H, fill_color=C_ACCENT3)

add_rect(s7, Inches(0.5), Inches(0.3), Inches(0.9), Inches(0.35), fill_color=C_ACCENT3)
add_text(s7, "05", Inches(0.5), Inches(0.3), Inches(0.9), Inches(0.35),
         font_size=Pt(14), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(s7, "预期效益", Inches(1.5), Inches(0.3), Inches(8), Inches(0.45),
         font_size=Pt(26), bold=True, color=C_WHITE)
add_rect(s7, Inches(0.5), Inches(0.82), Inches(12.3), Inches(0.03), fill_color=C_CARD)
add_text(s7, "社会价值 × 经济回报 × 数据资产 — 三维价值体系",
         Inches(0.5), Inches(0.95), Inches(12), Inches(0.4),
         font_size=Pt(14), color=C_GRAY)

# Big metrics row
metrics = [
    (C_ACCENT1, "万分之一", "编校差错率目标\nAI初筛+人工复核"),
    (C_ACCENT2, "35%", "出版周期缩短\nAI承担60%重复工作"),
    (C_GREEN,   "40%", "人才成才周期缩短\nAI导师赋能新编辑"),
    (C_ACCENT3, "10×", "交易效率提升\n意图驱动智能结算"),
]
for i, (accent, num, label) in enumerate(metrics):
    bx = Inches(0.4 + i * 3.2)
    by = Inches(1.5)
    add_rect(s7, bx, by, Inches(3.0), Inches(1.8),
             fill_color=C_CARD, line_color=accent, line_width=Pt(1.5))
    add_rect(s7, bx, by, Inches(3.0), Inches(0.07), fill_color=accent)
    add_text(s7, num, bx + Inches(0.1), by + Inches(0.1),
             Inches(2.8), Inches(0.9),
             font_size=Pt(36), bold=True, color=accent, align=PP_ALIGN.CENTER)
    add_text(s7, label, bx + Inches(0.1), by + Inches(1.0),
             Inches(2.8), Inches(0.7),
             font_size=Pt(11), color=C_LIGHT, align=PP_ALIGN.CENTER)

# Two columns: Social + Economic
# Social
add_rect(s7, Inches(0.4), Inches(3.5), Inches(6.0), Inches(3.7),
         fill_color=C_CARD, line_color=C_ACCENT1, line_width=Pt(1))
add_rect(s7, Inches(0.4), Inches(3.5), Inches(6.0), Inches(0.5), fill_color=C_ACCENT1)
add_text(s7, "社会效益", Inches(0.55), Inches(3.52),
         Inches(5.7), Inches(0.45),
         font_size=Pt(16), bold=True, color=C_BG_DARK)
social = [
    "✦  质量提升：差错率降至万分之一以下",
    "✦  人才培育：AI导师缩短成才周期40%",
    "✦  就业促进：降低门槛，赋能高校学生、退休教师",
    "✦  行业规范：AI定义「好质量」，标准化非标服务",
    "✦  资源协同：全省编校力量智能调度，削峰填谷",
]
add_text(s7, "\n".join(social), Inches(0.55), Inches(4.1),
         Inches(5.7), Inches(3.0),
         font_size=Pt(12), color=C_LIGHT)

# Economic
add_rect(s7, Inches(6.8), Inches(3.5), Inches(6.1), Inches(3.7),
         fill_color=C_CARD, line_color=C_ACCENT2, line_width=Pt(1))
add_rect(s7, Inches(6.8), Inches(3.5), Inches(6.1), Inches(0.5), fill_color=C_ACCENT2)
add_text(s7, "经济效益", Inches(6.95), Inches(3.52),
         Inches(5.8), Inches(0.45),
         font_size=Pt(16), bold=True, color=C_BG_DARK)
economic = [
    "✦  降本增效：AI承担60%重复工作，周期缩短35%",
    "✦  多元盈利：中介费 + AI工具订阅 + 增值服务",
    "✦  数据资产：汇聚高质量语料，构建垂直知识库",
    "✦  估值目标：不仅是5000万，目标亿级行业估值",
    "✦  规模效应：全国推广后，网络效应指数级增长",
]
add_text(s7, "\n".join(economic), Inches(6.95), Inches(4.1),
         Inches(5.8), Inches(3.0),
         font_size=Pt(12), color=C_LIGHT)

print("Slide 7 done.")

# ══════════════════════════════════════════════════════════════
# SLIDE 8 — INVESTMENT HIGHLIGHTS
# ══════════════════════════════════════════════════════════════
s8 = blank_slide(prs)
fill_bg(s8, C_BG_DARK)
add_rect(s8, Inches(0), Inches(0), Inches(0.08), H, fill_color=C_ACCENT2)

add_rect(s8, Inches(0.5), Inches(0.3), Inches(0.9), Inches(0.35), fill_color=C_ACCENT2)
add_text(s8, "06", Inches(0.5), Inches(0.3), Inches(0.9), Inches(0.35),
         font_size=Pt(14), bold=True, color=C_BG_DARK, align=PP_ALIGN.CENTER)
add_text(s8, "投资亮点", Inches(1.5), Inches(0.3), Inches(8), Inches(0.45),
         font_size=Pt(26), bold=True, color=C_WHITE)
add_rect(s8, Inches(0.5), Inches(0.82), Inches(12.3), Inches(0.03), fill_color=C_CARD)
add_text(s8, "为什么是现在？为什么是我们？",
         Inches(0.5), Inches(0.95), Inches(12), Inches(0.4),
         font_size=Pt(14), color=C_GRAY)

highlights = [
    (C_ACCENT1, "🎯", "蓝海赛道",
     "出版编校是万亿级文化产业的核心环节，AI 渗透率几乎为零，先发优势极为显著"),
    (C_ACCENT2, "🔬", "技术壁垒",
     "语义指纹算法 + 垂直知识库 + 数字信用体系，三重护城河，难以复制"),
    (C_GREEN, "📊", "数据飞轮",
     "每一次编校交互都在沉淀高质量语料，数据资产随规模指数级增长，形成强大壁垒"),
    (C_ACCENT3, "🏛️", "政策顺风",
     "国家大力推进出版业数字化转型，国产大模型政策支持，自主可控战略高度契合"),
    (C_ACCENT1, "🌐", "网络效应",
     "平台价值随用户增长而增长，出版社越多→编辑越多→数据越好→匹配越精准"),
    (C_ACCENT2, "💰", "多元变现",
     "SaaS订阅 + 交易佣金 + 数据服务 + API授权，多条盈利路径，抗风险能力强"),
]
for i, (accent, icon, title, desc) in enumerate(highlights):
    col = i % 3
    row = i // 3
    bx = Inches(0.4 + col * 4.3)
    by = Inches(1.4 + row * 2.7)
    bw, bh = Inches(4.0), Inches(2.4)
    add_rect(s8, bx, by, bw, bh, fill_color=C_CARD, line_color=accent, line_width=Pt(1.2))
    add_rect(s8, bx, by, bw, Inches(0.07), fill_color=accent)
    # Icon + title row
    add_text(s8, icon, bx + Inches(0.15), by + Inches(0.12),
             Inches(0.5), Inches(0.5), font_size=Pt(20), color=accent)
    add_text(s8, title, bx + Inches(0.7), by + Inches(0.15),
             Inches(3.1), Inches(0.45),
             font_size=Pt(17), bold=True, color=accent)
    add_rect(s8, bx + Inches(0.15), by + Inches(0.65), Inches(3.5), Inches(0.03), fill_color=accent)
    add_text(s8, desc, bx + Inches(0.15), by + Inches(0.75),
             Inches(3.7), Inches(1.55),
             font_size=Pt(12), color=C_LIGHT)

print("Slide 8 done.")

# ══════════════════════════════════════════════════════════════
# SLIDE 9 — CLOSING / CALL TO ACTION
# ══════════════════════════════════════════════════════════════
s9 = blank_slide(prs)
fill_bg(s9, C_BG_DARK)

# Full-width top accent
add_rect(s9, Inches(0), Inches(0), W, Inches(0.08), fill_color=C_ACCENT1)
add_rect(s9, Inches(0), H - Inches(0.08), W, Inches(0.08), fill_color=C_ACCENT1)

# Background decorative circles
add_circle(s9, Inches(1.5), Inches(5.5), Inches(5.0), RGBColor(0x0C, 0x14, 0x35))
add_circle(s9, Inches(11.5), Inches(2.0), Inches(4.0), RGBColor(0x0C, 0x14, 0x35))

# Central content
add_text(s9, "编得好", Inches(2.5), Inches(1.2), Inches(8.5), Inches(1.5),
         font_size=Pt(72), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(s9, "AI-Native 出版编校协作平台",
         Inches(2.5), Inches(2.7), Inches(8.5), Inches(0.6),
         font_size=Pt(24), color=C_ACCENT1, align=PP_ALIGN.CENTER)

# Divider
add_rect(s9, Inches(4.5), Inches(3.45), Inches(4.5), Inches(0.04), fill_color=C_ACCENT2)

# Tagline
add_text(s9, "不是修补，而是基因重组",
         Inches(2.5), Inches(3.6), Inches(8.5), Inches(0.6),
         font_size=Pt(20), color=C_ACCENT2, align=PP_ALIGN.CENTER, italic=True)

# Strategy pillars
pillars = ["应用先行", "自主可控", "数据为王", "敏捷迭代"]
for i, p in enumerate(pillars):
    bx = Inches(1.8 + i * 2.5)
    add_rect(s9, bx, Inches(4.4), Inches(2.1), Inches(0.65),
             fill_color=C_CARD, line_color=C_ACCENT1, line_width=Pt(1))
    add_text(s9, p, bx + Inches(0.1), Inches(4.42),
             Inches(1.9), Inches(0.6),
             font_size=Pt(15), bold=True, color=C_ACCENT1, align=PP_ALIGN.CENTER)

# CTA
add_rect(s9, Inches(3.5), Inches(5.3), Inches(6.5), Inches(0.9),
         fill_color=C_ACCENT1, line_color=None)
add_text(s9, "期待与您共同开创出版行业的 AI 新纪元",
         Inches(3.5), Inches(5.35), Inches(6.5), Inches(0.8),
         font_size=Pt(16), bold=True, color=C_BG_DARK, align=PP_ALIGN.CENTER)

# Footer
add_text(s9, "立足山东  ·  辐射全国  ·  引领出版行业从「人力密集型」向「数智协同型」跃迁",
         Inches(1.0), Inches(6.5), Inches(11.5), Inches(0.4),
         font_size=Pt(12), color=C_GRAY, align=PP_ALIGN.CENTER)

print("Slide 9 done.")

# ══════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════
out_path = "E:/yw/agiatme/goose/output/编得好_投资人PPT.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Total slides: {len(prs.slides)}")
