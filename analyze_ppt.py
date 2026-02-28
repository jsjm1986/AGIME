# -*- coding: utf-8 -*-
from pptx import Presentation

prs = Presentation('./data/workspaces/698616a1980c003c66f6421e/missions/40204152-de4f-41ac-bca3-17265a7f1112/documents/doc_b75f92060612ba80.pptx')
print(f'Slide size: {prs.slide_width} x {prs.slide_height}')

for i, slide in enumerate(prs.slides):
    layout_name = slide.slide_layout.name
    print(f'\n=== Slide {i+1}: layout="{layout_name}" ===')
    for shape in slide.shapes:
        txt = ""
        fs = ""
        if shape.has_text_frame:
            txt = shape.text_frame.text[:80].replace('\n', '|')
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    if r.font.size:
                        fs = f' sz={r.font.size}'
                    break
                break
        print(f'  [{shape.shape_type}] pos=({shape.left},{shape.top}) sz=({shape.width},{shape.height}){fs} "{txt}"')
