from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── 颜色主题 ──────────────────────────────────────────────
C_DARK   = RGBColor(0x1A, 0x23, 0x3A)   # 深海蓝（背景/标题栏）
C_ACCENT = RGBColor(0x00, 0x8B, 0xD8)   # 亮蓝（强调色）
C_LIGHT  = RGBColor(0xF0, 0xF4, 0xF8)   # 浅灰蓝（内容背景）
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_TEXT   = RGBColor(0x1A, 0x23, 0x3A)
C_GRAY   = RGBColor(0x55, 0x65, 0x7A)
C_TAG    = RGBColor(0xE8, 0xF4, 0xFD)   # 技能标签背景

W, H = Inches(13.33), Inches(7.5)       # 16:9 宽屏

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]            # 空白版式

# ── 工具函数 ──────────────────────────────────────────────
def add_rect(slide, l, t, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=C_TEXT,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = wrap
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.name  = "微软雅黑"
    return txb

def add_para(tf, text, size=14, bold=False, color=C_TEXT,
             align=PP_ALIGN.LEFT, space_before=0):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.name  = "微软雅黑"
    return p

def slide_header(slide, title, subtitle=None):
    """顶部深色标题栏"""
    add_rect(slide, 0, 0, W, Inches(1.1), fill=C_DARK)
    add_text(slide, title,
             Inches(0.4), Inches(0.15), Inches(10), Inches(0.7),
             size=28, bold=True, color=C_WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.4), Inches(0.72), Inches(10), Inches(0.35),
                 size=13, color=RGBColor(0xA0, 0xC8, 0xE8))
    # 底部装饰线
    add_rect(slide, 0, H - Inches(0.08), W, Inches(0.08), fill=C_ACCENT)

# ══════════════════════════════════════════════════════════
# Slide 1 — 封面
# ══════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
add_rect(s1, 0, 0, W, H, fill=C_DARK)
# 左侧亮色竖条
add_rect(s1, 0, 0, Inches(0.18), H, fill=C_ACCENT)
# 姓名
add_text(s1, "冯  瑞",
         Inches(0.6), Inches(1.8), Inches(8), Inches(1.4),
         size=64, bold=True, color=C_WHITE)
# 职位
add_text(s1, "AI Agent 开发工程师  ·  LLM 应用架构师  ·  AI 平台研发工程师",
         Inches(0.6), Inches(3.2), Inches(11), Inches(0.6),
         size=20, color=RGBColor(0xA0, 0xC8, 0xE8))
# 分隔线
add_rect(s1, Inches(0.6), Inches(3.9), Inches(5), Inches(0.04), fill=C_ACCENT)
# 联系信息
add_text(s1, "📱 18611101112     🐙 github.com/jsjm1986",
         Inches(0.6), Inches(4.1), Inches(9), Inches(0.5),
         size=15, color=RGBColor(0xCC, 0xDD, 0xEE))
# 右下角装饰
add_rect(s1, W - Inches(3.5), H - Inches(3.5), Inches(3.5), Inches(3.5),
         fill=RGBColor(0x22, 0x30, 0x50))
add_text(s1, "AI\nAgent",
         W - Inches(3.2), H - Inches(3.0), Inches(3.0), Inches(2.5),
         size=72, bold=True, color=RGBColor(0x2A, 0x3D, 0x60), align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════
# Slide 2 — 个人简介
# ══════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
add_rect(s2, 0, 0, W, H, fill=C_LIGHT)
slide_header(s2, "个人简介", "About Me")

# 简介文字框
add_rect(s2, Inches(0.5), Inches(1.3), Inches(12.3), Inches(2.5), fill=C_WHITE, line=C_ACCENT)
txb = slide_header  # reuse nothing, just add textbox
tb2 = s2.shapes.add_textbox(Inches(0.7), Inches(1.45), Inches(11.9), Inches(2.2))
tb2.word_wrap = True
tf2 = tb2.text_frame
tf2.word_wrap = True
p0 = tf2.paragraphs[0]
p0.alignment = PP_ALIGN.LEFT
r0 = p0.add_run()
r0.text = ("专注于 AI Agent 系统架构与工程落地的开发者，具备多个企业级 AI 平台的完整设计与开发经验。"
           "深入实践 LangGraph 多 Agent 协作、smolagents 框架、工作流引擎等核心技术，"
           "在 Agent 架构设计、容错机制、状态管理、性能优化等方面积累了丰富的工程经验。"
           "积极参与开源社区，完成 Block 公司 Goose AI Agent 项目的中文本地化。"
           "具备从系统架构设计到生产部署的全流程交付能力。")
r0.font.size = Pt(15)
r0.font.color.rgb = C_TEXT
r0.font.name = "微软雅黑"

# 三个亮点卡片
cards = [
    ("🏗️", "架构设计", "LangGraph / smolagents\n多 Agent 协作系统"),
    ("⚙️", "工程落地", "从设计到生产部署\n全流程交付能力"),
    ("🌐", "开源贡献", "Goose AI Agent\n中文本地化"),
]
for i, (icon, title, desc) in enumerate(cards):
    x = Inches(0.5 + i * 4.3)
    add_rect(s2, x, Inches(4.1), Inches(4.0), Inches(2.8), fill=C_WHITE, line=C_ACCENT)
    add_text(s2, icon, x + Inches(0.2), Inches(4.2), Inches(0.8), Inches(0.7), size=28)
    add_text(s2, title, x + Inches(0.2), Inches(4.95), Inches(3.5), Inches(0.45),
             size=16, bold=True, color=C_ACCENT)
    add_text(s2, desc, x + Inches(0.2), Inches(5.45), Inches(3.5), Inches(1.2),
             size=13, color=C_GRAY)

# ══════════════════════════════════════════════════════════
# Slide 3 — 核心能力
# ══════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
add_rect(s3, 0, 0, W, H, fill=C_LIGHT)
slide_header(s3, "核心能力", "Core Skills")

skills = [
    ("🤖 Agent 系统架构",
     "LangGraph / smolagents 多 Agent 协作\nSupervisor Pattern · ReAct 推理\nHuman-in-the-Loop · 状态机 · DAG 工作流"),
    ("🧠 LLM 工程化",
     "Prompt Engineering · DSPy 自进化\nFew-shot / CoT · 语义缓存\n向量检索 · Token 成本优化"),
    ("⚙️ 后端开发",
     "Python · FastAPI · Node.js\nSQLAlchemy · Pydantic V2\nWebSocket · Redis · Docker"),
    ("🗄️ 数据层",
     "MySQL · PostgreSQL · SQLite\nRedis · LanceDB / Chroma 向量库\nText2SQL · Schema 自动提取"),
    ("🖥️ 前端开发",
     "React 18 · TypeScript\nAnt Design · React Flow · ECharts\nZustand 状态管理"),
    ("🔧 工程实践",
     "Git 工作流 · CI/CD\nDocker 容器化 · Gunicorn/Uvicorn\n日志监控 · 单元测试"),
]

cols, rows = 3, 2
cw, ch = Inches(4.1), Inches(2.5)
for idx, (title, body) in enumerate(skills):
    col = idx % cols
    row = idx // cols
    x = Inches(0.4 + col * 4.3)
    y = Inches(1.25 + row * 2.7)
    add_rect(s3, x, y, cw, ch, fill=C_WHITE, line=C_ACCENT)
    add_rect(s3, x, y, cw, Inches(0.38), fill=C_ACCENT)
    add_text(s3, title, x + Inches(0.12), y + Inches(0.04), cw - Inches(0.2), Inches(0.32),
             size=13, bold=True, color=C_WHITE)
    add_text(s3, body, x + Inches(0.12), y + Inches(0.45), cw - Inches(0.2), ch - Inches(0.55),
             size=11.5, color=C_TEXT)

# ══════════════════════════════════════════════════════════
# 通用项目页辅助函数
# ══════════════════════════════════════════════════════════
def project_slide(title, subtitle, role_stack, bullets, tag_color=C_ACCENT):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, W, H, fill=C_LIGHT)
    slide_header(s, title, subtitle)

    # 角色 & 技术栈标签行
    add_rect(s, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.42), fill=C_WHITE, line=tag_color)
    add_text(s, role_stack,
             Inches(0.65), Inches(1.25), Inches(12.0), Inches(0.35),
             size=11, color=C_GRAY)

    # 要点列表
    tb = s.shapes.add_textbox(Inches(0.5), Inches(1.75), Inches(12.3), Inches(5.4))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for bullet in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        r = p.add_run()
        r.text = bullet
        r.font.size = Pt(13.5)
        r.font.color.rgb = C_TEXT
        r.font.name = "微软雅黑"
    return s

# ══════════════════════════════════════════════════════════
# Slide 4 — 项目：DataPilot
# ══════════════════════════════════════════════════════════
project_slide(
    "DataPilot — 企业级 Agentic BI 平台",
    "Project Experience 1",
    "角色：独立设计与开发  |  技术栈：Python · LangGraph 1.x · FastAPI · React 18 · TypeScript · DeepSeek API · Redis · LanceDB · WebSocket",
    [
        "🏗️  【系统架构】基于 LangGraph 1.x 设计 Supervisor Pattern 多 Agent 协作架构，实现自然语言→数据洞察端到端自动化，含 6 个专业 Agent 分工协作",
        "🤖  【Agent 矩阵】AmbiResolver(歧义消解) · DataSniper(Schema 侦察) · LogicArchitect(SQL 生成) · Judge(执行校验) · VizExpert(可视化) · CacheCheck(缓存检查)",
        "🔄  【Human-in-the-Loop】基于 LangGraph interrupt() 机制实现人机协作，支持 5 种歧义类型（时间/指标/范围/粒度/实体）的智能识别与主动澄清",
        "🛡️  【容错设计】三级自愈策略：Agent 级重试(3次) → Supervisor ReAct 推理修复 → Human Handoff 断点续执行",
        "⚡  【智能路由】三级复杂度路由引擎：简单查询串行执行 · 中等查询并行处理 · 复杂查询启用 ReAct 多步推理",
        "📈  【自进化系统】集成 DSPy 框架，自动收集成功案例优化 Prompt，支持模块版本管理与性能回滚",
        "💾  【缓存架构】三层语义缓存（内存 L1 → Redis L2 → LanceDB L3），双重匹配策略（精确哈希 + 向量相似度 ≥ 0.85）",
        "🔒  【安全机制】SQLGlot AST 解析防注入 · E2B 沙箱隔离代码执行 · EXPLAIN 成本熔断 · JWT 认证",
    ]
)

# ══════════════════════════════════════════════════════════
# Slide 5 — 项目：Multi-Agent System
# ══════════════════════════════════════════════════════════
project_slide(
    "Multi-Agent System — 基于 smolagents 的多智能体平台",
    "Project Experience 2",
    "角色：独立设计与开发  |  技术栈：Python 3.12 · smolagents · FastAPI · Uvicorn/Gunicorn · SQLAlchemy · Docker · Pydantic V2",
    [
        "🏗️  【架构设计】基于 HuggingFace smolagents 框架构建模块化多智能体系统，采用 Manager-Worker 模式实现任务调度与分发",
        "🤖  【Agent 矩阵】Manager Agent(任务调度) · SQL Agent(数据库操作) · Code Agent(代码执行) · Text Agent(文本处理)，职责单一、接口清晰",
        "🔧  【工具系统】search_tool(语义搜索) · webpage_tool(网页抓取转 Markdown) · code_tool(安全沙箱执行) · sql_tool(SQLAlchemy 查询)",
        "🧠  【多步推理】支持 Thought→Code→Observation 迭代机制处理复杂任务，同时支持 single_step 模式优化简单任务",
        "🚀  【生产部署】Docker 容器化部署，Gunicorn + Uvicorn 生产配置，支持 Nginx 反向代理与负载均衡",
        "⚙️  【配置管理】Pydantic V2 类型安全配置，自动创建目录结构，完善的日志系统与错误处理",
    ]
)

# ══════════════════════════════════════════════════════════
# Slide 6 — 项目：LLM2Workflow
# ══════════════════════════════════════════════════════════
project_slide(
    "LLM2Workflow — 自然语言工作流设计平台",
    "Project Experience 3",
    "角色：独立设计与开发  |  技术栈：React · TypeScript · React Flow · Ant Design · FastAPI · LangChain · PostgreSQL · DeepSeek API",
    [
        "💬  【核心功能】通过自然语言描述自动生成可视化工作流，支持流式输出实时展示生成过程",
        "🔲  【节点系统】13 种节点类型：HTTP 请求 · 数据转换 · LLM 调用 · 数据库操作 · 文件处理 · 邮件 · 图像处理 · 调度器 · WebSocket · 消息队列 · 循环 · 并行 · 聚合",
        "🖱️  【可视化编辑】基于 React Flow 实现拖拽式节点编排，智能连线与自动布局，双击节点配置参数",
        "⚙️  【工作流引擎】支持条件分支、循环控制、并行执行、结果聚合等复杂业务逻辑编排",
        "📊  【执行监控】实时执行状态追踪，完善的异常处理与错误提示机制",
    ]
)

# ══════════════════════════════════════════════════════════
# Slide 7 — 其他项目
# ══════════════════════════════════════════════════════════
project_slide(
    "其他项目经验",
    "Project Experience 4",
    "Goose AI Agent 中文本地化 · 领域 AI Agent（医疗/心理咨询）",
    [
        "🐦  【Goose AI Agent — 开源贡献】",
        "     Block 公司（Square/CashApp 母公司）开源本地 AI 代理，完成核心文档简体中文翻译与本地化",
        "     深入理解 MCP(Model Context Protocol) 服务器集成、多 LLM Provider 支持、Extension 扩展机制等架构设计",
        "     推动 AI Agent 工具在中文开发者社区的普及与应用",
        "",
        "🏥  【领域 AI Agent — 医疗/心理咨询】",
        "     中医问诊 Agent：完整问诊→辨证→治疗方案流程，7 种辨证维度综合分析，结构化诊疗报告生成",
        "     心理咨询 Agent：五阶段状态机管理会话流程，实时情绪识别与响应策略调整，危机识别安全机制",
        "     工程特点：流式响应 · 语音输入 · 响应式多端适配 · 完善数据导出",
    ]
)

# ══════════════════════════════════════════════════════════
# Slide 8 — 工程能力亮点
# ══════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(BLANK)
add_rect(s8, 0, 0, W, H, fill=C_LIGHT)
slide_header(s8, "工程能力亮点", "Engineering Highlights")

highlights = [
    ("🏗️", "架构设计",
     "熟练运用 LangGraph、smolagents 等框架\n掌握 Supervisor、ReAct、Human-in-the-Loop 等 Agent 设计模式"),
    ("🛡️", "系统可靠性",
     "注重容错机制设计\n多级重试、自愈修复、断点续执行、优雅降级"),
    ("⚡", "性能优化",
     "多层缓存架构、智能路由、并行执行\nToken 成本控制等性能优化实践"),
    ("🚀", "生产部署",
     "Docker 容器化、Gunicorn 生产配置\n日志监控、安全防护等 DevOps 实践"),
    ("🌐", "开源贡献",
     "积极参与 AI Agent 开源社区\n具备良好的技术文档能力和协作意识"),
]

for i, (icon, title, desc) in enumerate(highlights):
    col = i % 3
    row = i // 3
    x = Inches(0.4 + col * 4.3)
    y = Inches(1.3 + row * 2.85)
    bw = Inches(4.0) if i < 3 else Inches(6.15)
    if i == 3:
        x = Inches(0.4)
    elif i == 4:
        x = Inches(6.9)
    add_rect(s8, x, y, bw, Inches(2.6), fill=C_WHITE, line=C_ACCENT)
    add_rect(s8, x, y, Inches(0.5), Inches(2.6), fill=C_ACCENT)
    add_text(s8, icon, x + Inches(0.08), y + Inches(0.85), Inches(0.4), Inches(0.6), size=20, color=C_WHITE)
    add_text(s8, title, x + Inches(0.6), y + Inches(0.15), bw - Inches(0.7), Inches(0.45),
             size=15, bold=True, color=C_ACCENT)
    add_text(s8, desc, x + Inches(0.6), y + Inches(0.65), bw - Inches(0.7), Inches(1.8),
             size=12.5, color=C_TEXT)

# ══════════════════════════════════════════════════════════
# Slide 9 — 自我评价 & 求职意向
# ══════════════════════════════════════════════════════════
s9 = prs.slides.add_slide(BLANK)
add_rect(s9, 0, 0, W, H, fill=C_DARK)
add_rect(s9, 0, 0, Inches(0.18), H, fill=C_ACCENT)
slide_header(s9, "自我评价 & 求职意向", "Self-Assessment & Career Goal")

evals = [
    "🔬  对 AI Agent 技术有深入理解，持续跟进 LangGraph、AutoGPT、CrewAI 等前沿框架发展",
    "🏆  注重工程质量，强调系统可靠性、可维护性、可扩展性，具备生产级系统设计能力",
    "🚀  具备优秀的自驱力，能够独立完成从架构设计到部署上线的全流程交付",
    "💡  善于将复杂 AI 能力工程化为稳定可靠的产品功能，兼具技术深度与产品思维",
]

tb9 = s9.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(12.0), Inches(3.5))
tb9.word_wrap = True
tf9 = tb9.text_frame
tf9.word_wrap = True
for i, ev in enumerate(evals):
    p = tf9.paragraphs[0] if i == 0 else tf9.add_paragraph()
    p.space_before = Pt(8)
    r = p.add_run()
    r.text = ev
    r.font.size = Pt(15)
    r.font.color.rgb = RGBColor(0xDD, 0xEE, 0xFF)
    r.font.name = "微软雅黑"

# 求职意向标签
add_rect(s9, Inches(0.6), Inches(5.2), Inches(12.0), Inches(0.05), fill=C_ACCENT)
add_text(s9, "求职意向",
         Inches(0.6), Inches(5.35), Inches(3), Inches(0.4),
         size=14, bold=True, color=C_ACCENT)
targets = ["AI Agent 开发工程师", "LLM 应用架构师", "AI 平台研发工程师"]
for i, t in enumerate(targets):
    bx = Inches(0.6 + i * 4.1)
    add_rect(s9, bx, Inches(5.85), Inches(3.8), Inches(0.7),
             fill=RGBColor(0x00, 0x5A, 0x8E), line=C_ACCENT)
    add_text(s9, t, bx + Inches(0.15), Inches(5.9), Inches(3.5), Inches(0.55),
             size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════
# 保存
# ══════════════════════════════════════════════════════════
out_path = os.path.join(os.path.dirname(__file__), "resume_冯瑞.pptx")
prs.save(out_path)
print(f"✅ 已生成 {len(prs.slides)} 页 PPT：{out_path}")
