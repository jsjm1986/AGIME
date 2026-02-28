from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colors ──────────────────────────────────────────────
BG       = RGBColor(0x0A, 0x0E, 0x1A)
CARD     = RGBColor(0x14, 0x19, 0x29)
BLUE     = RGBColor(0x1E, 0x6F, 0xFF)
GOLD     = RGBColor(0xF5, 0xC8, 0x42)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
SUBTEXT  = RGBColor(0xA8, 0xB8, 0xD8)
SIDEBAR  = RGBColor(0x0D, 0x12, 0x22)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

def rgb(r,g,b): return RGBColor(r,g,b)

def add_rect(slide, x, y, w, h, fill, alpha=None):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.line.fill.background()
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    return s

def add_text(slide, text, x, y, w, h, size, bold=False, color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "微软雅黑"
    return tb

def bg(slide):
    add_rect(slide, 0, 0, W, H, BG)

def sidebar(slide, label, num):
    add_rect(slide, 0, 0, Inches(0.08), H, BLUE)
    add_rect(slide, Inches(0.08), 0, Inches(1.5), H, SIDEBAR)
    add_text(slide, num, Inches(0.12), Inches(0.2), Inches(1.3), Inches(0.5), 11, color=GOLD, align=PP_ALIGN.CENTER)
    add_text(slide, label, Inches(0.12), Inches(0.8), Inches(1.3), Inches(5), 11, color=SUBTEXT, align=PP_ALIGN.CENTER)

def content_title(slide, title, subtitle=None):
    add_text(slide, title, Inches(1.9), Inches(0.35), Inches(10.5), Inches(0.7), 30, bold=True, color=WHITE)
    # gold underline
    add_rect(slide, Inches(1.9), Inches(1.05), Inches(1.2), Inches(0.04), GOLD)
    if subtitle:
        add_text(slide, subtitle, Inches(1.9), Inches(1.15), Inches(10), Inches(0.4), 14, color=SUBTEXT)

def card(slide, x, y, w, h, title, body, icon=""):
    add_rect(slide, x, y, w, h, CARD)
    # top accent line
    add_rect(slide, x, y, w, Inches(0.05), BLUE)
    if icon:
        add_text(slide, icon, x+Inches(0.2), y+Inches(0.12), Inches(0.5), Inches(0.5), 22, color=GOLD)
    add_text(slide, title, x+Inches(0.2), y+Inches(0.15), w-Inches(0.4), Inches(0.45), 15, bold=True, color=WHITE)
    add_text(slide, body, x+Inches(0.2), y+Inches(0.65), w-Inches(0.4), h-Inches(0.8), 12, color=SUBTEXT)

# ── Slide 01: Cover ──────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
# gradient overlay rect
add_rect(sl, 0, 0, W, H, rgb(0x0A,0x0E,0x1A))
# blue glow bar left
add_rect(sl, 0, 0, Inches(0.12), H, BLUE)
# decorative horizontal lines
for i, y_pos in enumerate([Inches(2.8), Inches(2.88), Inches(2.96)]):
    w_line = Inches(13.33 - i*1.5)
    add_rect(sl, 0, y_pos, w_line, Inches(0.015), rgb(0x1E,0x6F,0xFF) if i==0 else rgb(0x14,0x19,0x29))

add_text(sl, "编得好", Inches(1.2), Inches(1.0), Inches(11), Inches(1.5), 72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "全国首个 AI 原生出版编校协作平台", Inches(1.2), Inches(2.5), Inches(11), Inches(0.7), 24, color=GOLD, align=PP_ALIGN.CENTER)
add_rect(sl, Inches(5.2), Inches(3.1), Inches(2.9), Inches(0.05), BLUE)
add_text(sl, "AI 经纪人  ·  AI 助理  ·  AI 审计员", Inches(1.2), Inches(3.3), Inches(11), Inches(0.5), 16, color=SUBTEXT, align=PP_ALIGN.CENTER)
add_text(sl, "项目立项建议书  |  投资人版", Inches(1.2), Inches(6.5), Inches(11), Inches(0.5), 13, color=SUBTEXT, align=PP_ALIGN.CENTER)

# ── Slide 02: Agenda ─────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl); sidebar(sl, "目录\nAGENDA", "02")
content_title(sl, "议程概览", "六大核心板块，带您全面了解编得好")
items = [
    ("01", "行业痛点", "出版行业三大困局"),
    ("02", "解决方案", "AI原生重构生产关系"),
    ("03", "核心功能", "四大智能模块全景"),
    ("04", "实施路线", "12个月三阶段计划"),
    ("05", "商业模式", "多元盈利矩阵"),
    ("06", "预期效益", "社会与经济双重价值"),
]
for i, (num, title, sub) in enumerate(items):
    col = i % 3; row = i // 3
    x = Inches(1.9) + col * Inches(3.7)
    y = Inches(1.6) + row * Inches(2.3)
    add_rect(sl, x, y, Inches(3.4), Inches(2.0), CARD)
    add_rect(sl, x, y, Inches(3.4), Inches(0.05), BLUE)
    add_text(sl, num, x+Inches(0.2), y+Inches(0.1), Inches(0.6), Inches(0.5), 28, bold=True, color=BLUE)
    add_text(sl, title, x+Inches(0.2), y+Inches(0.65), Inches(3.0), Inches(0.45), 16, bold=True, color=WHITE)
    add_text(sl, sub, x+Inches(0.2), y+Inches(1.15), Inches(3.0), Inches(0.6), 12, color=SUBTEXT)

# ── Slide 03: Pain Points ────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl); sidebar(sl, "行业痛点\nPAIN POINTS", "03")
content_title(sl, "出版行业的三大困局", "传统模式已无法满足数字化时代的需求")
pains = [
    ("资源分布不均", "供需信息不对称", "编校资源集中于少数城市，优质稿件与顶尖编辑难以跨地域匹配。传统平台依赖人工搜索，效率低下，资源孤岛严重。"),
    ("效率低下", "流程繁琐耗时", "依赖QQ/微信发包，缺乏工具支撑。合同、发票、个税手动处理，一个项目的沟通成本占总工时30%以上。"),
    ("质量难评", "信用体系空白", "评价主观性强，无客观标准。劣质编辑与优质编辑同价竞争，劣币驱逐良币，行业信任危机严重。"),
]
for i, (title, sub, body) in enumerate(pains):
    x = Inches(1.9) + i * Inches(3.7)
    add_rect(sl, x, Inches(1.55), Inches(3.4), Inches(5.5), CARD)
    add_rect(sl, x, Inches(1.55), Inches(3.4), Inches(0.06), GOLD)
    num_str = ["01", "02", "03"][i]
    add_text(sl, num_str, x+Inches(0.2), Inches(1.65), Inches(0.8), Inches(0.6), 32, bold=True, color=GOLD)
    add_text(sl, title, x+Inches(0.2), Inches(2.35), Inches(3.0), Inches(0.45), 16, bold=True, color=WHITE)
    add_text(sl, sub, x+Inches(0.2), Inches(2.82), Inches(3.0), Inches(0.35), 12, bold=False, color=BLUE)
    add_text(sl, body, x+Inches(0.2), Inches(3.25), Inches(3.0), Inches(3.5), 11, color=SUBTEXT)

# ── Slide 04: Solution ───────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl); sidebar(sl, "解决方案\nSOLUTION", "04")
content_title(sl, "1+3+N 架构：AI 原生重构生产关系", "一个平台 · 三类智能体 · N个生态伙伴")

# Center hub
cx, cy = Inches(6.5), Inches(4.0)
add_rect(sl, cx-Inches(1.1), cy-Inches(0.7), Inches(2.2), Inches(1.4), BLUE)
add_text(sl, "编得好平台", cx-Inches(1.0), cy-Inches(0.55), Inches(2.0), Inches(0.4), 14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "AI 原生协作中枢", cx-Inches(1.0), cy-Inches(0.1), Inches(2.0), Inches(0.35), 11, color=GOLD, align=PP_ALIGN.CENTER)

agents = [
    (Inches(2.0), Inches(2.0), "AI 经纪人", "智能匹配\n供需双方"),
    (Inches(2.0), Inches(5.2), "AI 助理", "全程辅助\n编校工作"),
    (Inches(10.0), Inches(3.6), "AI 审计员", "质量评估\n信用背书"),
]
for ax, ay, aname, adesc in agents:
    add_rect(sl, ax, ay, Inches(2.0), Inches(1.2), CARD)
    add_rect(sl, ax, ay, Inches(2.0), Inches(0.05), GOLD)
    add_text(sl, aname, ax+Inches(0.15), ay+Inches(0.1), Inches(1.7), Inches(0.4), 14, bold=True, color=GOLD)
    add_text(sl, adesc, ax+Inches(0.15), ay+Inches(0.55), Inches(1.7), Inches(0.55), 11, color=SUBTEXT)

# ── Slide 05: Core Features ──────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl); sidebar(sl, "核心功能\nFEATURES", "05")
content_title(sl, "四大智能模块全景", "覆盖出版全流程的 AI 能力矩阵")
feats = [
    ("智能匹配引擎", "基于NLP语义分析，精准匹配稿件与编辑专长。\n• 多维度画像建模\n• 实时供需动态平衡\n• 跨地域资源调度"),
    ("AI 辅助编校", "大模型驱动的编校助手，提升效率60%以上。\n• 错别字/语法自动检测\n• 风格一致性分析\n• 智能批注与建议"),
    ("全流程管理", "从接单到交付的数字化闭环管理。\n• 合同自动生成\n• 进度实时追踪\n• 发票/个税一键处理"),
    ("信用评估体系", "客观量化的编辑能力评估与信用积累。\n• AI质量评分\n• 历史数据建模\n• 动态信用等级"),
]
for i, (title, body) in enumerate(feats):
    col = i % 2; row = i // 2
    x = Inches(1.9) + col * Inches(5.6)
    y = Inches(1.55) + row * Inches(2.7)
    add_rect(sl, x, y, Inches(5.2), Inches(2.45), CARD)
    add_rect(sl, x, y, Inches(5.2), Inches(0.06), BLUE)
    add_text(sl, ["01","02","03","04"][i], x+Inches(0.2), y+Inches(0.1), Inches(0.6), Inches(0.45), 22, bold=True, color=BLUE)
    add_text(sl, title, x+Inches(0.85), y+Inches(0.1), Inches(4.1), Inches(0.45), 15, bold=True, color=WHITE)
    add_text(sl, body, x+Inches(0.2), y+Inches(0.65), Inches(4.8), Inches(1.65), 11, color=SUBTEXT)

# ── Slide 06: Roadmap ────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl); sidebar(sl, "实施路线\nROADMAP", "06")
content_title(sl, "12个月三阶段实施计划", "稳步推进，快速验证，规模扩张")
phases = [
    ("Phase 1", "0-4个月", "MVP验证期",
     "• 核心匹配功能上线\n• 种子用户招募（50家出版社）\n• 基础AI编校工具发布\n• 完成天使轮融资"),
    ("Phase 2", "5-8个月", "规模增长期",
     "• AI审计员模块上线\n• 用户规模突破500家\n• 月GMV达到500万元\n• 完成Pre-A轮融资"),
    ("Phase 3", "9-12个月", "生态构建期",
     "• 开放平台API\n• 引入N个生态合作伙伴\n• 月GMV突破2000万元\n• 启动A轮融资"),
]
for i, (phase, period, title, body) in enumerate(phases):
    x = Inches(1.9) + i * Inches(3.7)
    add_rect(sl, x, Inches(1.55), Inches(3.4), Inches(5.5), CARD)
    # Phase color bar
    colors = [RGBColor(0x1A,0x6B,0xC8), RGBColor(0xC8,0x9A,0x1A), RGBColor(0x1A,0xC8,0x7A)]
    add_rect(sl, x, Inches(1.55), Inches(3.4), Inches(0.06), colors[i])
    add_text(sl, phase, x+Inches(0.2), Inches(1.65), Inches(1.5), Inches(0.4), 13, bold=True, color=colors[i])
    add_text(sl, period, x+Inches(1.8), Inches(1.65), Inches(1.4), Inches(0.4), 11, color=SUBTEXT)
    add_text(sl, title, x+Inches(0.2), Inches(2.15), Inches(3.0), Inches(0.45), 15, bold=True, color=WHITE)
    add_text(sl, body, x+Inches(0.2), Inches(2.7), Inches(3.0), Inches(4.0), 11, color=SUBTEXT)

# ── Slide 07: Business Model ─────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl); sidebar(sl, "商业模式\nBUSINESS", "07")
content_title(sl, "多元收入结构，高粘性平台生态", "平台抽佣 + SaaS订阅 + 增值服务三驾马车")
models = [
    ("平台交易抽佣", "主营收入", "对每笔编校交易收取8-12%服务费。\n预计占总收入60%，随GMV线性增长。\n\n目标：Year1 GMV 5000万元"),
    ("SaaS工具订阅", "稳定现金流", "出版社订阅AI编校工具套件。\n基础版免费，专业版¥2980/月/席位。\n\n目标：Year1 付费席位1000+"),
    ("增值服务", "高毛利收入", "• 优先推荐排名\n• 信用认证服务\n• 数据分析报告\n• 企业定制培训\n\n毛利率 > 80%"),
]
for i, (title, sub, body) in enumerate(models):
    x = Inches(1.9) + i * Inches(3.7)
    add_rect(sl, x, Inches(1.55), Inches(3.4), Inches(5.5), CARD)
    add_rect(sl, x, Inches(1.55), Inches(3.4), Inches(0.06), GOLD)
    add_text(sl, title, x+Inches(0.2), Inches(1.7), Inches(3.0), Inches(0.45), 15, bold=True, color=WHITE)
    add_text(sl, sub, x+Inches(0.2), Inches(2.2), Inches(3.0), Inches(0.35), 12, color=GOLD)
    add_text(sl, body, x+Inches(0.2), Inches(2.65), Inches(3.0), Inches(4.0), 11, color=SUBTEXT)

# ── Slide 08: Market Size ────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl); sidebar(sl, "市场规模\nMARKET", "08")
content_title(sl, "千亿级蓝海市场，AI重塑行业格局", "中国出版编辑服务市场规模持续扩大")
stats = [
    ("¥1200亿", "中国出版业年产值", "全国580余家出版社，年出版图书50万种"),
    ("80万+", "注册编辑从业者", "其中自由编辑占比超40%，数字化需求迫切"),
    ("60%", "效率提升空间", "传统流程痛点显著，AI替代潜力巨大"),
    ("¥50亿", "可寻址市场(SAM)", "编辑服务数字化平台的直接目标市场"),
]
for i, (num, title, desc) in enumerate(stats):
    col = i % 2; row = i // 2
    x = Inches(1.9) + col * Inches(5.6)
    y = Inches(1.55) + row * Inches(2.7)
    add_rect(sl, x, y, Inches(5.2), Inches(2.45), CARD)
    add_text(sl, num, x+Inches(0.3), y+Inches(0.2), Inches(4.6), Inches(0.9), 36, bold=True, color=GOLD)
    add_text(sl, title, x+Inches(0.3), y+Inches(1.1), Inches(4.6), Inches(0.4), 14, bold=True, color=WHITE)
    add_text(sl, desc, x+Inches(0.3), y+Inches(1.55), Inches(4.6), Inches(0.75), 11, color=SUBTEXT)

# ── Slide 09: Financing Ask ──────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl); sidebar(sl, "融资计划\nFINANCING", "09")
content_title(sl, "寻求天使轮融资 ¥1500万", "加速产品研发与市场拓展")
add_rect(sl, Inches(1.9), Inches(1.55), Inches(11.1), Inches(1.1), CARD)
add_rect(sl, Inches(1.9), Inches(1.55), Inches(11.1), Inches(0.06), GOLD)
add_text(sl, "融资金额：¥1500万元  ·  出让股权：10%  ·  投后估值：¥1.5亿元", Inches(2.1), Inches(1.65), Inches(10.7), Inches(0.85), 16, bold=True, color=GOLD)
uses = [
    ("40%  ·  ¥600万", "产品研发", "AI引擎迭代、移动端开发、安全架构升级"),
    ("30%  ·  ¥450万", "市场推广", "品牌建设、渠道拓展、种子用户运营"),
    ("20%  ·  ¥300万", "团队扩张", "核心技术与运营人才招募"),
    ("10%  ·  ¥150万", "运营储备", "日常运营、法务合规、应急储备"),
]
for i, (pct, title, desc) in enumerate(uses):
    x = Inches(1.9) + i * Inches(2.8)
    add_rect(sl, x, Inches(2.85), Inches(2.6), Inches(3.8), CARD)
    add_rect(sl, x, Inches(2.85), Inches(2.6), Inches(0.06), BLUE)
    add_text(sl, pct, x+Inches(0.2), Inches(2.95), Inches(2.2), Inches(0.5), 14, bold=True, color=BLUE)
    add_text(sl, title, x+Inches(0.2), Inches(3.5), Inches(2.2), Inches(0.4), 13, bold=True, color=WHITE)
    add_text(sl, desc, x+Inches(0.2), Inches(3.95), Inches(2.2), Inches(0.6), 10, color=SUBTEXT)

# ── Slide 10: Closing ────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
# Full-width gold accent line
add_rect(sl, Inches(0), Inches(3.5), Inches(13.33), Inches(0.04), GOLD)
add_text(sl, "编得好", Inches(1.5), Inches(1.2), Inches(10), Inches(1.2), 60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "让每一本好书，都能遇见最好的编辑", Inches(1.5), Inches(2.5), Inches(10), Inches(0.7), 22, color=SUBTEXT, align=PP_ALIGN.CENTER)
add_rect(sl, Inches(4.5), Inches(3.65), Inches(4.33), Inches(0.04), RGBColor(0x33,0x33,0x55))
add_text(sl, "期待与您携手，共同重塑出版行业的未来", Inches(1.5), Inches(3.8), Inches(10), Inches(0.6), 16, color=GOLD, align=PP_ALIGN.CENTER)
add_text(sl, "联系我们  ·  contact@biandehao.com  ·  www.biandehao.com", Inches(1.5), Inches(5.5), Inches(10), Inches(0.5), 13, color=SUBTEXT, align=PP_ALIGN.CENTER)
add_text(sl, "© 2025 编得好科技有限公司  ·  保密文件，仅供投资参考", Inches(1.5), Inches(6.5), Inches(10), Inches(0.4), 10, color=RGBColor(0x44,0x44,0x66), align=PP_ALIGN.CENTER)

# ── Save ─────────────────────────────────────────────────
import os
os.makedirs("output", exist_ok=True)
out = "output/编得好_投资人版.pptx"
prs.save(out)
print(f"✅ Saved: {out}  ({len(prs.slides)} slides)")
